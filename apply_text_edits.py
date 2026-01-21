#!/usr/bin/env python3

import argparse
import os
import tempfile

# PIP3 modules
import pptx
import yaml

# local repo modules
import slide_deck_pipeline.csv_schema as csv_schema
import slide_deck_pipeline.pptx_hash as pptx_hash
import slide_deck_pipeline.pptx_text as pptx_text
import slide_deck_pipeline.text_boxes as text_boxes
import slide_deck_pipeline.soffice_tools as soffice_tools


MAX_BULLET_DEPTH = 4


#============================================
def parse_args() -> argparse.Namespace:
	"""
	Parse command-line arguments.
	"""
	parser = argparse.ArgumentParser(
		description="Apply text edits from a YAML patch file."
	)
	parser.add_argument(
		"-i",
		"--input",
		dest="input_path",
		required=True,
		help="Input PPTX or ODP file",
	)
	parser.add_argument(
		"-p",
		"--patches",
		dest="patch_path",
		required=True,
		help="YAML patch file",
	)
	parser.add_argument(
		"-o",
		"--output",
		dest="output_path",
		default="",
		help="Output PPTX or ODP path (default: <input>_edited.pptx)",
	)
	parser.add_argument(
		"-f",
		"--force",
		dest="force",
		help="Apply edits even if text hashes mismatch",
		action="store_true",
	)
	parser.add_argument(
		"-F",
		"--no-force",
		dest="force",
		help="Skip edits when text hashes mismatch",
		action="store_false",
	)
	parser.set_defaults(force=False)
	parser.add_argument(
		"-s",
		"--include-subtitle",
		dest="include_subtitle",
		help="Include subtitle placeholders in matching",
		action="store_true",
	)
	parser.add_argument(
		"-S",
		"--no-include-subtitle",
		dest="include_subtitle",
		help="Skip subtitle placeholders in matching",
		action="store_false",
	)
	parser.set_defaults(include_subtitle=False)
	parser.add_argument(
		"-r",
		"--include-footer",
		dest="include_footer",
		help="Include footer placeholders in matching",
		action="store_true",
	)
	parser.add_argument(
		"-R",
		"--no-include-footer",
		dest="include_footer",
		help="Skip footer placeholders in matching",
		action="store_false",
	)
	parser.set_defaults(include_footer=False)
	args = parser.parse_args()
	return args


#============================================
#============================================
def resolve_input_pptx(input_path: str, temp_dir: str | None) -> tuple[str, str]:
	"""
	Return a PPTX path and a source basename for matching.

	Args:
		input_path: Input PPTX or ODP path.
		temp_dir: Temporary directory for conversions, or None.

	Returns:
		tuple[str, str]: Resolved PPTX path and source basename.
	"""
	source_name = os.path.basename(input_path)
	lowered = input_path.lower()
	if lowered.endswith(".pptx"):
		return (input_path, source_name)
	if lowered.endswith(".odp"):
		if not temp_dir:
			raise ValueError("Temporary directory required for ODP conversion.")
		pptx_path = soffice_tools.convert_odp_to_pptx(input_path, temp_dir)
		return (pptx_path, source_name)
	raise ValueError("Input must be a .pptx or .odp file.")


#============================================
def parse_text_lines(text_value: str) -> list[tuple[int, str]]:
	"""
	Parse text into indentation levels.

	Args:
		text_value: Text with leading tabs for indentation.

	Returns:
		list[tuple[int, str]]: List of (level, text).
	"""
	if not text_value:
		return []
	lines = []
	cleaned = text_value.replace("\r\n", "\n").replace("\r", "\n")
	for raw_line in cleaned.split("\n"):
		if raw_line == "":
			lines.append((0, ""))
			continue
		level = len(raw_line) - len(raw_line.lstrip("\t"))
		text = raw_line.lstrip("\t")
		lines.append((level, text))
	return lines


#============================================
def render_bullets(items, level: int = 0) -> list[str]:
	"""
	Render nested bullet lists into tab-indented lines.

	Args:
		items: Bullet items.
		level: Indentation level.

	Returns:
		list[str]: Rendered lines.
	"""
	lines = []
	if items is None:
		return lines
	if level >= MAX_BULLET_DEPTH:
		raise ValueError("Bullet nesting depth exceeds the maximum.")
	if isinstance(items, str):
		lines.append("\t" * level + items)
		return lines
	if not isinstance(items, list):
		raise ValueError("Bullets must be a list or string.")
	for item in items:
		if isinstance(item, str):
			lines.append("\t" * level + item)
			continue
		if isinstance(item, list):
			lines.extend(render_bullets(item, level + 1))
			continue
		raise ValueError("Bullets must contain strings or lists only.")
	return lines


#============================================
def resolve_box_text(box: dict[str, object]) -> str:
	"""
	Resolve new text content for a box.

	Args:
		box: Box entry from YAML.

	Returns:
		str: Text content.
	"""
	if "bullets" in box and box["bullets"] is not None:
		lines = render_bullets(box["bullets"])
		return "\n".join(lines)
	return str(box.get("text", ""))


#============================================
def should_skip_box(box: dict[str, object]) -> bool:
	"""
	Determine if a box should be skipped.

	Args:
		box: Box entry from YAML.

	Returns:
		bool: True if the box is locked.
	"""
	if box.get("locked"):
		return True
	ed_status = str(box.get("edit_status", "")).strip().lower()
	if ed_status in ("locked", "skip", "frozen"):
		return True
	return False


#============================================
def set_shape_text(shape, text_value: str) -> None:
	"""
	Set text on a shape with indentation levels.

	Args:
		shape: Shape instance.
		text_value: Text content.
	"""
	if not getattr(shape, "has_text_frame", False):
		return
	text_frame = shape.text_frame
	text_frame.clear()
	lines = parse_text_lines(text_value)
	if not lines:
		return
	for index, (level, text) in enumerate(lines):
		if index == 0:
			paragraph = text_frame.paragraphs[0]
		else:
			paragraph = text_frame.add_paragraph()
		paragraph.text = text
		paragraph.level = level


#============================================
def apply_text_edits(
	input_path: str,
	patch_path: str,
	output_path: str,
	force: bool,
	include_subtitle: bool,
	include_footer: bool,
) -> None:
	"""
	Apply text edits to a deck.

	Args:
		input_path: Input PPTX or ODP path.
		patch_path: YAML patch file.
		output_path: Output PPTX or ODP path.
		force: Apply edits even if text hashes mismatch.
		include_subtitle: Include subtitle placeholders in matching.
		include_footer: Include footer placeholders in matching.
	"""
	with open(patch_path, "r", encoding="utf-8") as handle:
		payload = yaml.safe_load(handle)
	if not isinstance(payload, dict):
		raise ValueError("Patch file must be a YAML mapping.")
	version = payload.get("version")
	if str(version).isdigit():
		version_value = int(version)
	else:
		raise ValueError("Unsupported patch version.")
	if version_value != 1:
		raise ValueError("Unsupported patch version.")
	patches = payload.get("patches", [])
	if not isinstance(patches, list):
		raise ValueError("Patch file must contain a patches list.")
	source_name = str(payload.get("source_pptx", "")).strip()
	input_base = os.path.basename(input_path)
	if source_name and source_name != input_base:
		print("WARN: patch source_pptx does not match input file basename.")

	needs_conversion = input_path.lower().endswith(".odp")
	if needs_conversion:
		with tempfile.TemporaryDirectory() as temp_dir:
			pptx_path, _ = resolve_input_pptx(input_path, temp_dir)
			apply_and_save(
				pptx_path,
				patches,
				output_path,
				force,
				include_subtitle,
				include_footer,
			)
		return
	pptx_path, _ = resolve_input_pptx(input_path, None)
	apply_and_save(
		pptx_path,
		patches,
		output_path,
		force,
		include_subtitle,
		include_footer,
	)


#============================================
def apply_and_save(
	pptx_path: str,
	patches: list[dict[str, object]],
	output_path: str,
	force: bool,
	include_subtitle: bool,
	include_footer: bool,
) -> None:
	"""
	Apply edits to a PPTX and save to output.

	Args:
		pptx_path: Input PPTX path.
		patches: Patch entries.
		output_path: Output PPTX or ODP path.
		force: Apply edits even if text hashes mismatch.
		include_subtitle: Include subtitle placeholders in matching.
		include_footer: Include footer placeholders in matching.
	"""
	presentation = pptx.Presentation(pptx_path)
	updated = 0
	skipped = 0
	missing = 0
	mismatched = 0
	slide_hash_mismatch = 0
	for patch in patches:
		if not isinstance(patch, dict):
			continue
		slide_index = patch.get("source_slide_index")
		if slide_index is None:
			missing += 1
			continue
		slide_index_text = str(slide_index)
		if not slide_index_text.isdigit():
			missing += 1
			continue
		slide_number = int(slide_index_text)
		if slide_number < 1 or slide_number > len(presentation.slides):
			missing += 1
			continue
		slide = presentation.slides[slide_number - 1]
		notes_text = pptx_text.extract_notes_text(slide)
		current_hash, _, _ = pptx_hash.compute_slide_hash_from_slide(
			slide,
			notes_text,
		)
		expected_hash = str(patch.get("slide_hash", ""))
		if not expected_hash or expected_hash != current_hash:
			slide_hash_mismatch += 1
			continue
		boxes, _ = text_boxes.collect_text_boxes(
			slide,
			include_subtitle,
			include_footer,
			include_fallback=True,
		)
		box_map = {box["box_id"]: box for box in boxes}
		boxes_data = patch.get("boxes", [])
		if not isinstance(boxes_data, list):
			missing += 1
			continue
		for box in boxes_data:
			if not isinstance(box, dict):
				continue
			if should_skip_box(box):
				skipped += 1
				continue
			box_id = str(box.get("box_id", ""))
			if not box_id:
				missing += 1
				continue
			if box_id == "notes":
				new_text = resolve_box_text(box)
				if not force:
					current_notes_hash = csv_schema.compute_text_hash(notes_text)
					expected_notes_hash = str(box.get("text_hash_before", ""))
					if expected_notes_hash and expected_notes_hash != current_notes_hash:
						mismatched += 1
						continue
				set_notes_text(slide, new_text)
				updated += 1
				continue
			box_meta = box_map.get(box_id)
			if not box_meta:
				missing += 1
				continue
			shape = box_meta["shape"]
			current_text = text_boxes.extract_text_block(shape)
			if not force:
				expected_hash = str(box.get("text_hash_before", ""))
				current_hash = csv_schema.compute_text_hash(current_text)
				if expected_hash and expected_hash != current_hash:
					mismatched += 1
					continue
			new_text = resolve_box_text(box)
			set_shape_text(shape, new_text)
			updated += 1

	if output_path.lower().endswith(".odp"):
		with tempfile.TemporaryDirectory() as temp_dir:
			temp_pptx = os.path.join(temp_dir, "edited.pptx")
			presentation.save(temp_pptx)
			soffice_tools.convert_pptx_to_odp(temp_pptx, output_path)
	else:
		presentation.save(output_path)

	print(f"Updated blocks: {updated}")
	print(f"Skipped locked blocks: {skipped}")
	print(f"Missing targets: {missing}")
	print(f"Text hash mismatches: {mismatched}")
	print(f"Slide hash mismatches: {slide_hash_mismatch}")


#============================================
def set_notes_text(slide: pptx.slide.Slide, notes_text: str) -> None:
	"""
	Set speaker notes text on a slide.

	Args:
		slide: Slide instance.
		notes_text: Notes text.
	"""
	if not notes_text and notes_text != "":
		return
	notes_slide = slide.notes_slide
	if not notes_slide:
		return
	notes_frame = notes_slide.notes_text_frame
	if not notes_frame:
		return
	notes_frame.text = notes_text


#============================================
def main() -> None:
	"""
	Main entry point.
	"""
	args = parse_args()
	output_path = args.output_path
	if not output_path:
		base_name = os.path.splitext(args.input_path)[0]
		output_path = f"{base_name}_edited.pptx"
	apply_text_edits(
		args.input_path,
		args.patch_path,
		output_path,
		args.force,
		args.include_subtitle,
		args.include_footer,
	)


if __name__ == "__main__":
	main()
