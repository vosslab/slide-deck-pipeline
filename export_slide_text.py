#!/usr/bin/env python3

import argparse
import os
import tempfile

# PIP3 modules
import pptx
import yaml

# local repo modules
import slide_deck_pipeline.csv_schema as csv_schema
import slide_deck_pipeline.pptx_text as pptx_text
import slide_deck_pipeline.text_boxes as text_boxes
import slide_deck_pipeline.soffice_tools as soffice_tools


#============================================
def parse_args() -> argparse.Namespace:
	"""
	Parse command-line arguments.
	"""
	parser = argparse.ArgumentParser(
		description="Export slide text blocks to a YAML patch file."
	)
	parser.add_argument(
		"-i",
		"--input",
		dest="input_path",
		required=True,
		help="Input PPTX or ODP file",
	)
	parser.add_argument(
		"-o",
		"--output",
		dest="output_path",
		default="",
		help="Output YAML path (default: <input>_text_edits.yaml)",
	)
	parser.add_argument(
		"-n",
		"--include-notes",
		dest="include_notes",
		help="Include speaker notes blocks",
		action="store_true",
	)
	parser.add_argument(
		"-N",
		"--no-include-notes",
		dest="include_notes",
		help="Skip speaker notes blocks",
		action="store_false",
	)
	parser.set_defaults(include_notes=False)
	parser.add_argument(
		"-s",
		"--include-subtitle",
		dest="include_subtitle",
		help="Include subtitle placeholders",
		action="store_true",
	)
	parser.add_argument(
		"-S",
		"--no-include-subtitle",
		dest="include_subtitle",
		help="Skip subtitle placeholders",
		action="store_false",
	)
	parser.set_defaults(include_subtitle=False)
	parser.add_argument(
		"-f",
		"--include-footer",
		dest="include_footer",
		help="Include footer placeholders",
		action="store_true",
	)
	parser.add_argument(
		"-F",
		"--no-include-footer",
		dest="include_footer",
		help="Skip footer placeholders",
		action="store_false",
	)
	parser.set_defaults(include_footer=False)
	args = parser.parse_args()
	return args


#============================================
#============================================
def resolve_input_pptx(input_path: str, temp_dir: str | None) -> tuple[str, str]:
	"""
	Return a PPTX path and a source basename for YAML output.

	Args:
		input_path: Input PPTX or ODP path.
		temp_dir: Temporary directory for conversions, or None.

	Returns:
		tuple[str, str]: Resolved PPTX path and source basename.
	"""
	source_name = os.path.basename(input_path)
	lowered = input_path.lower()
	if lowered.endswith(".pptx"):
		return (input_path, source_name)
	if lowered.endswith(".odp"):
		if not temp_dir:
			raise ValueError("Temporary directory required for ODP conversion.")
		pptx_path = soffice_tools.convert_odp_to_pptx(input_path, temp_dir)
		return (pptx_path, source_name)
	raise ValueError("Input must be a .pptx or .odp file.")


#============================================
def build_box_record(shape, box_meta: dict[str, object]) -> dict[str, str]:
	"""
	Build a YAML box record from a shape.

	Args:
		shape: Shape instance.
		box_meta: Metadata for the box.

	Returns:
		dict[str, str]: Box record.
	"""
	text_value = text_boxes.extract_text_block(shape)
	box_record = {
		"box_id": box_meta["box_id"],
		"text_hash_before": csv_schema.compute_text_hash(text_value),
		"text": text_value,
	}
	shape_name = box_meta.get("shape_name", "")
	if shape_name:
		box_record["shape_name"] = shape_name
	placeholder_type = box_meta.get("placeholder_type", "")
	if placeholder_type:
		box_record["placeholder_type"] = placeholder_type
	return box_record


#============================================
def export_slide_text(
	input_path: str,
	output_path: str,
	include_notes: bool,
	include_subtitle: bool,
	include_footer: bool,
) -> None:
	"""
	Export slide text blocks to YAML.

	Args:
		input_path: Input PPTX or ODP path.
		output_path: Output YAML path.
		include_notes: Include speaker notes blocks.
		include_subtitle: Include subtitle placeholders.
		include_footer: Include footer placeholders.
	"""
	needs_conversion = input_path.lower().endswith(".odp")
	if needs_conversion:
		with tempfile.TemporaryDirectory() as temp_dir:
			pptx_path, source_name = resolve_input_pptx(input_path, temp_dir)
			write_yaml(
				pptx_path,
				source_name,
				output_path,
				include_notes,
				include_subtitle,
				include_footer,
			)
		return
	pptx_path, source_name = resolve_input_pptx(input_path, None)
	write_yaml(
		pptx_path,
		source_name,
		output_path,
		include_notes,
		include_subtitle,
		include_footer,
	)


#============================================
def write_yaml(
	pptx_path: str,
	source_name: str,
	output_path: str,
	include_notes: bool,
	include_subtitle: bool,
	include_footer: bool,
) -> None:
	"""
	Write YAML from a PPTX path.

	Args:
		pptx_path: PPTX path.
		source_name: Source deck basename.
		output_path: Output YAML path.
		include_notes: Include speaker notes blocks.
		include_subtitle: Include subtitle placeholders.
		include_footer: Include footer placeholders.
	"""
	presentation = pptx.Presentation(pptx_path)
	patches = []
	fallback_slides = []
	box_count = 0
	for index, slide in enumerate(presentation.slides, 1):
		slide_text = pptx_text.extract_slide_text(slide)
		notes_text = pptx_text.extract_notes_text(slide)
		slide_hash = csv_schema.compute_slide_hash(slide_text, notes_text)
		boxes, used_fallback = text_boxes.collect_text_boxes(
			slide,
			include_subtitle,
			include_footer,
			include_fallback=True,
		)
		if used_fallback:
			fallback_slides.append(index)
		box_records = []
		for box_meta in boxes:
			shape = box_meta["shape"]
			box_records.append(build_box_record(shape, box_meta))
		if include_notes:
			box_records.append(
				{
					"box_id": "notes",
					"text_hash_before": csv_schema.compute_text_hash(notes_text),
					"text": notes_text,
					"placeholder_type": "notes",
				}
			)
		if not box_records:
			continue
		box_count += len(box_records)
		patches.append(
			{
				"source_slide_index": index,
				"slide_hash": slide_hash,
				"boxes": box_records,
			}
		)
	payload = {
		"version": 1,
		"source_pptx": source_name,
		"patches": patches,
	}
	with open(output_path, "w", encoding="utf-8") as handle:
		yaml.safe_dump(
			payload,
			handle,
			sort_keys=False,
			default_flow_style=False,
			allow_unicode=False,
		)
	print(f"Exported {len(patches)} slides with {box_count} text blocks.")
	if fallback_slides:
		listed = ", ".join(str(idx) for idx in fallback_slides)
		print(f"Fallback shape matching used on slides: {listed}")


#============================================
def main() -> None:
	"""
	Main entry point.
	"""
	args = parse_args()
	output_path = args.output_path
	if not output_path:
		base_name = os.path.splitext(args.input_path)[0]
		output_path = f"{base_name}_text_edits.yaml"
	export_slide_text(
		args.input_path,
		output_path,
		args.include_notes,
		args.include_subtitle,
		args.include_footer,
	)


if __name__ == "__main__":
	main()
