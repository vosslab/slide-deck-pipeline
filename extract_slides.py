import argparse
import hashlib
import os
import shutil
import subprocess
import tempfile

# PIP3 modules
import pptx
import pptx.enum.shapes

# local repo modules
import slide_csv


LAYOUT_HINT_MAP = {
	"title_and_content": "title_and_content",
	"title_and_object": "title_and_content",
	"title_only": "section_header",
	"section_header": "section_header",
	"two_content": "two_column",
	"two_column": "two_column",
	"blank": "blank",
}


#============================================
def parse_args() -> argparse.Namespace:
	"""
	Parse command-line arguments.
	"""
	parser = argparse.ArgumentParser(
		description="Extract slide content and assets into a CSV."
	)
	parser.add_argument(
		"-i",
		"--input",
		dest="input_path",
		required=True,
		help="Input PPTX or ODP file",
	)
	parser.add_argument(
		"-o",
		"--output",
		dest="output_csv",
		required=True,
		help="Output CSV path",
	)
	parser.add_argument(
		"-a",
		"--assets-dir",
		dest="assets_dir",
		default="",
		help="Assets directory (default: <output_csv>_assets)",
	)
	args = parser.parse_args()
	return args


#============================================
def convert_odp_to_pptx(odp_path: str, work_dir: str) -> str:
	"""
	Convert an ODP file to PPTX using soffice.

	Args:
		odp_path: Path to the ODP file.
		work_dir: Output directory for the converted PPTX.

	Returns:
		str: Path to the converted PPTX file.
	"""
	soffice_bin = shutil.which("soffice")
	if not soffice_bin:
		raise FileNotFoundError("soffice not found. Install LibreOffice to convert ODP.")
	command = [
		soffice_bin,
		"--headless",
		"--convert-to",
		"pptx",
		"--outdir",
		work_dir,
		odp_path,
	]
	result = subprocess.run(command, capture_output=True, text=True, cwd=work_dir)
	if result.returncode != 0:
		message = result.stderr.strip() or result.stdout.strip()
		raise RuntimeError(f"ODP conversion failed: {message}")
	base_name = os.path.splitext(os.path.basename(odp_path))[0]
	pptx_path = os.path.join(work_dir, f"{base_name}.pptx")
	if not os.path.exists(pptx_path):
		raise FileNotFoundError(f"Converted PPTX not found: {pptx_path}")
	return pptx_path


#============================================
def resolve_input_pptx(input_path: str, temp_dir: str | None) -> tuple[str, str]:
	"""
	Return a PPTX path and a source basename for CSV rows.

	Args:
		input_path: Input PPTX or ODP path.
		temp_dir: Temporary directory for conversions, or None.

	Returns:
		tuple[str, str]: Resolved PPTX path and source basename.
	"""
	source_name = os.path.basename(input_path)
	lowered = input_path.lower()
	if lowered.endswith(".pptx"):
		return (input_path, source_name)
	if lowered.endswith(".odp"):
		if not temp_dir:
			raise ValueError("Temporary directory required for ODP conversion.")
		pptx_path = convert_odp_to_pptx(input_path, temp_dir)
		return (pptx_path, source_name)
	raise ValueError("Input must be a .pptx or .odp file.")


#============================================
def normalize_layout_hint(layout_name: str) -> str:
	"""
	Normalize a PowerPoint layout name to a hint value.

	Args:
		layout_name: Layout name from pptx.

	Returns:
		str: Normalized layout hint.
	"""
	if not layout_name:
		return "custom"
	normalized = layout_name.strip().lower().replace(" ", "_")
	return LAYOUT_HINT_MAP.get(normalized, normalized)


#============================================
def extract_paragraph_lines(text_frame: pptx.text.text.TextFrame) -> list[str]:
	"""
	Extract text frame paragraphs with indentation markers.

	Args:
		text_frame: TextFrame instance.

	Returns:
		list[str]: Lines with leading tab indentation.
	"""
	lines = []
	for paragraph in text_frame.paragraphs:
		text = paragraph.text.strip()
		if not text:
			continue
		indent = "\t" * paragraph.level
		lines.append(f"{indent}{text}")
	return lines


#============================================
def extract_body_text(slide: pptx.slide.Slide) -> str:
	"""
	Extract body text from non-title text frames.

	Args:
		slide: Slide instance.

	Returns:
		str: Body text with newline separators.
	"""
	lines = []
	title_shape = slide.shapes.title
	for shape in slide.shapes:
		if title_shape and shape == title_shape:
			continue
		if not shape.has_text_frame:
			continue
		lines.extend(extract_paragraph_lines(shape.text_frame))
	return "\n".join(lines)


#============================================
def extract_notes_text(slide: pptx.slide.Slide) -> str:
	"""
	Extract speaker notes text from a slide.

	Args:
		slide: Slide instance.

	Returns:
		str: Notes text.
	"""
	if not slide.has_notes_slide:
		return ""
	notes_frame = slide.notes_slide.notes_text_frame
	if not notes_frame:
		return ""
	return notes_frame.text or ""


#============================================
def collect_slide_images(slide: pptx.slide.Slide) -> list[dict[str, str | bytes]]:
	"""
	Collect slide image blobs and hashes.

	Args:
		slide: Slide instance.

	Returns:
		list[dict[str, str | bytes]]: Image records with blob, ext, and hash.
	"""
	images = []
	for shape in slide.shapes:
		if shape.shape_type != pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
			continue
		blob = shape.image.blob
		ext = shape.image.ext
		digest = hashlib.sha256(blob).hexdigest()
		images.append({"blob": blob, "ext": ext, "hash": digest})
	return images


#============================================
def export_images(
	images: list[dict[str, str | bytes]],
	assets_dir: str,
	slide_uid: str,
) -> tuple[list[str], list[str]]:
	"""
	Export collected images to disk.

	Args:
		images: Image records with blob, ext, and hash.
		assets_dir: Asset output directory.
		slide_uid: Stable slide UID.

	Returns:
		tuple[list[str], list[str]]: Image refs and image hashes.
	"""
	image_refs = []
	image_hashes = []
	for index, image in enumerate(images, 1):
		blob = image["blob"]
		ext = image["ext"]
		digest = image["hash"]
		filename = f"{slide_uid}_img{index:02d}.{ext}"
		output_path = os.path.join(assets_dir, filename)
		with open(output_path, "wb") as handle:
			handle.write(blob)
		image_refs.append(filename)
		image_hashes.append(digest)
	return (image_refs, image_hashes)


#============================================
def build_slide_row(
	source_name: str,
	slide_index: int,
	title_text: str,
	body_text: str,
	notes_text: str,
	layout_hint: str,
	image_refs: list[str],
	image_hashes: list[str],
) -> dict[str, str]:
	"""
	Build a CSV row for a slide.

	Args:
		source_name: Source PPTX basename.
		slide_index: 1-based slide index.
		title_text: Title text.
		body_text: Body text.
		notes_text: Notes text.
		layout_hint: Layout hint.
		image_refs: Image refs list.
		image_hashes: Image hash list.

	Returns:
		dict[str, str]: CSV row.
	"""
	text_hash = slide_csv.compute_text_hash(title_text, body_text, notes_text)
	slide_fingerprint = slide_csv.compute_slide_fingerprint(
		title_text,
		body_text,
		notes_text,
		image_hashes,
	)
	slide_uid = slide_csv.compute_slide_uid(
		source_name,
		slide_index,
		title_text,
		body_text,
		notes_text,
		image_hashes,
	)
	return {
		"source_pptx": source_name,
		"source_slide_index": str(slide_index),
		"slide_uid": slide_uid,
		"title_text": title_text,
		"body_text": body_text,
		"notes_text": notes_text,
		"layout_hint": layout_hint,
		"image_refs": slide_csv.join_list_field(image_refs),
		"image_hashes": slide_csv.join_list_field(image_hashes),
		"text_hash": text_hash,
		"slide_fingerprint": slide_fingerprint,
	}


#============================================
def ensure_assets_dir(output_csv: str, assets_dir: str) -> str:
	"""
	Resolve and create the assets directory.

	Args:
		output_csv: Output CSV path.
		assets_dir: Assets directory or empty string.

	Returns:
		str: Resolved assets directory.
	"""
	resolved_dir = assets_dir
	if not resolved_dir:
		base_name = os.path.splitext(output_csv)[0]
		resolved_dir = f"{base_name}_assets"
	if not os.path.exists(resolved_dir):
		os.makedirs(resolved_dir, exist_ok=True)
	return resolved_dir


#============================================
def extract_slides_to_csv(input_path: str, output_csv: str, assets_dir: str) -> None:
	"""
	Extract slides to CSV and export images.

	Args:
		input_path: Input PPTX or ODP path.
		output_csv: Output CSV path.
		assets_dir: Assets directory.
	"""
	resolved_assets_dir = ensure_assets_dir(output_csv, assets_dir)
	needs_conversion = input_path.lower().endswith(".odp")
	if needs_conversion:
		with tempfile.TemporaryDirectory() as temp_dir:
			pptx_path, source_name = resolve_input_pptx(input_path, temp_dir)
			rows = extract_rows(pptx_path, source_name, resolved_assets_dir)
	else:
		pptx_path, source_name = resolve_input_pptx(input_path, None)
		rows = extract_rows(pptx_path, source_name, resolved_assets_dir)
	slide_csv.write_slide_csv(output_csv, rows)


#============================================
def extract_rows(
	pptx_path: str,
	source_name: str,
	assets_dir: str,
) -> list[dict[str, str]]:
	"""
	Extract rows from a PPTX path.

	Args:
		pptx_path: Path to PPTX.
		source_name: Source basename for CSV rows.
		assets_dir: Asset output directory.

	Returns:
		list[dict[str, str]]: CSV rows.
	"""
	presentation = pptx.Presentation(pptx_path)
	rows = []
	for index, slide in enumerate(presentation.slides, 1):
		title_text = ""
		if slide.shapes.title and slide.shapes.title.text_frame:
			title_text = slide.shapes.title.text_frame.text or ""
		body_text = extract_body_text(slide)
		notes_text = extract_notes_text(slide)
		layout_hint = normalize_layout_hint(slide.slide_layout.name)
		images = collect_slide_images(slide)
		image_hashes = [image["hash"] for image in images]
		image_refs = []
		row = build_slide_row(
			source_name,
			index,
			title_text,
			body_text,
			notes_text,
			layout_hint,
			image_refs,
			image_hashes,
		)
		if images:
			image_refs, image_hashes = export_images(
				images,
				assets_dir,
				row["slide_uid"],
			)
			row["image_refs"] = slide_csv.join_list_field(image_refs)
			row["image_hashes"] = slide_csv.join_list_field(image_hashes)
		rows.append(row)
	return rows


#============================================
def main() -> None:
	"""
	Main entry point.
	"""
	args = parse_args()
	extract_slides_to_csv(args.input_path, args.output_csv, args.assets_dir)


if __name__ == "__main__":
	main()
