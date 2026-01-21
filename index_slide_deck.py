#!/usr/bin/env python3

import argparse
import os
import tempfile

# PIP3 modules
import pptx
import pptx.enum.shapes

# local repo modules
import slide_deck_pipeline.pptx_hash as pptx_hash
import slide_deck_pipeline.pptx_text as pptx_text
import slide_deck_pipeline.csv_schema as csv_schema
import slide_deck_pipeline.soffice_tools as soffice_tools
import slide_deck_pipeline.layout_classifier as layout_classifier


#============================================
def parse_args() -> argparse.Namespace:
	"""
	Parse command-line arguments.
	"""
	parser = argparse.ArgumentParser(
		description="Index slide content into a CSV."
	)
	parser.add_argument(
		"-i",
		"--input",
		dest="input_path",
		required=True,
		help="Input PPTX or ODP file",
	)
	parser.add_argument(
		"-o",
		"--output",
		dest="output_csv",
		default="",
		help="Output CSV path (default: <input>.csv)",
	)
	args = parser.parse_args()
	return args


#============================================
#============================================
def resolve_input_pptx(input_path: str, temp_dir: str | None) -> tuple[str, str]:
	"""
	Return a PPTX path and a source basename for CSV rows.

	Args:
		input_path: Input PPTX or ODP path.
		temp_dir: Temporary directory for conversions, or None.

	Returns:
		tuple[str, str]: Resolved PPTX path and source basename.
	"""
	source_name = os.path.basename(input_path)
	lowered = input_path.lower()
	if lowered.endswith(".pptx"):
		return (input_path, source_name)
	if lowered.endswith(".odp"):
		if not temp_dir:
			raise ValueError("Temporary directory required for ODP conversion.")
		pptx_path = soffice_tools.convert_odp_to_pptx(input_path, temp_dir)
		return (pptx_path, source_name)
	raise ValueError("Input must be a .pptx or .odp file.")


#============================================
def extract_paragraph_lines(text_frame: pptx.text.text.TextFrame) -> list[str]:
	"""
	Extract text frame paragraphs with indentation markers.

	Args:
		text_frame: TextFrame instance.

	Returns:
		list[str]: Lines with leading tab indentation.
	"""
	return pptx_text.extract_paragraph_lines(text_frame)


#============================================
def extract_body_text(slide: pptx.slide.Slide) -> str:
	"""
	Extract body text from non-title text frames.

	Args:
		slide: Slide instance.

	Returns:
		str: Body text with newline separators.
	"""
	return pptx_text.extract_body_text(slide)


#============================================
def extract_notes_text(slide: pptx.slide.Slide) -> str:
	"""
	Extract speaker notes text from a slide.

	Args:
		slide: Slide instance.

	Returns:
		str: Notes text.
	"""
	return pptx_text.extract_notes_text(slide)


#============================================
def collect_asset_types(slide: pptx.slide.Slide) -> str:
	"""
	Summarize slide asset types for context.

	Args:
		slide: Slide instance.

	Returns:
		str: Asset type summary string.
	"""
	media_type = getattr(pptx.enum.shapes.MSO_SHAPE_TYPE, "MEDIA", None)
	counts = {
		"images": 0,
		"tables": 0,
		"charts": 0,
		"media": 0,
	}

	def scan_shape(shape) -> None:
		if (
			getattr(shape, "shape_type", None)
			== pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP
			and hasattr(shape, "shapes")
		):
			for nested in shape.shapes:
				scan_shape(nested)
			return
		if getattr(shape, "shape_type", None) == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
			counts["images"] += 1
		if getattr(shape, "has_table", False):
			counts["tables"] += 1
		if getattr(shape, "has_chart", False):
			counts["charts"] += 1
		if media_type and getattr(shape, "shape_type", None) == media_type:
			counts["media"] += 1

	for shape in slide.shapes:
		scan_shape(shape)

	labels = []
	if counts["images"] == 1:
		labels.append("image")
	elif counts["images"] > 1:
		labels.append(f"images_{counts['images']}")
	if counts["tables"]:
		labels.append("table")
	if counts["charts"]:
		labels.append("chart")
	if counts["media"]:
		labels.append("media")
	return "|".join(labels)


#============================================
def describe_shape_type(shape) -> str:
	"""
	Return a readable shape type name.

	Args:
		shape: Shape instance.

	Returns:
		str: Shape type name.
	"""
	shape_type = getattr(shape, "shape_type", None)
	if shape_type is None:
		return "unknown"
	return getattr(shape_type, "name", str(shape_type))


#============================================
def is_supported_shape(shape) -> bool:
	"""
	Check whether a shape is supported for indexing.

	Args:
		shape: Shape instance.

	Returns:
		bool: True if supported.
	"""
	if getattr(shape, "is_placeholder", False):
		return True
	if getattr(shape, "shape_type", None) == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
		return True
	if getattr(shape, "has_text_frame", False):
		return True
	if getattr(shape, "has_table", False):
		return True
	if getattr(shape, "has_chart", False):
		return True
	return False


#============================================
def collect_unsupported_shapes(slide: pptx.slide.Slide) -> list[str]:
	"""
	Collect unsupported shape types from a slide.

	Args:
		slide: Slide instance.

	Returns:
		list[str]: Shape type names.
	"""
	unsupported = []
	for shape in slide.shapes:
		if (
			getattr(shape, "shape_type", None)
			== pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP
			and hasattr(shape, "shapes")
		):
			for nested in shape.shapes:
				if not is_supported_shape(nested):
					unsupported.append(describe_shape_type(nested))
			continue
		if not is_supported_shape(shape):
			unsupported.append(describe_shape_type(shape))
	return unsupported


#============================================
def resolve_master_name(slide: pptx.slide.Slide) -> tuple[str, str | None]:
	"""
	Resolve master name for a slide.

	Args:
		slide: Slide instance.

	Returns:
		tuple[str, str | None]: master name, warning text.
	"""
	try:
		layout = slide.slide_layout
	except Exception as exc:
		message = f"layout lookup failed: {exc}"
		return ("custom", message)
	master_name = "custom"
	slide_master = getattr(layout, "slide_master", None)
	if slide_master and getattr(slide_master, "name", ""):
		master_name = slide_master.name
	return (master_name, None)


#============================================
def report_index_warnings(
	unsupported_shapes: dict[int, list[str]],
	layout_errors: dict[int, str],
) -> None:
	"""
	Print a simple indexing warning report.

	Args:
		unsupported_shapes: Slide indices with unsupported shapes.
	"""
	if not unsupported_shapes and not layout_errors:
		return
	print("Index warnings:")
	if unsupported_shapes:
		for slide_index, shapes in sorted(unsupported_shapes.items()):
			joined = ", ".join(shapes)
			print(f"- Slide {slide_index} unsupported shapes: {joined}")
	if layout_errors:
		for slide_index, message in sorted(layout_errors.items()):
			print(f"- Slide {slide_index} layout warning: {message}")


#============================================
def build_slide_row(
	source_name: str,
	slide_index: int,
	title_text: str,
	body_text: str,
	notes_text: str,
	slide_hash: str,
	master_name: str,
	layout_type: str,
	asset_types: str,
) -> dict[str, str]:
	"""
	Build a CSV row for a slide.

	Args:
		source_name: Source PPTX basename.
		slide_index: 1-based slide index.
		title_text: Title text.
		body_text: Body text.
	notes_text: Notes text.
	slide_hash: Slide hash.
	master_name: Template master name.
	layout_type: Computed semantic layout type.
	asset_types: Asset type summary.

	Returns:
		dict[str, str]: CSV row.
	"""
	return {
		"source_pptx": source_name,
		"source_slide_index": str(slide_index),
		"slide_hash": slide_hash,
		"master_name": master_name,
		"layout_type": layout_type,
		"asset_types": asset_types,
		"title_text": title_text,
		"body_text": body_text,
		"notes_text": notes_text,
	}


#============================================
def index_slides_to_csv(input_path: str, output_csv: str) -> None:
	"""
	Index slides to CSV.

	Args:
		input_path: Input PPTX or ODP path.
		output_csv: Output CSV path.
	"""
	needs_conversion = input_path.lower().endswith(".odp")
	if needs_conversion:
		with tempfile.TemporaryDirectory() as temp_dir:
			pptx_path, source_name = resolve_input_pptx(input_path, temp_dir)
			rows = index_rows(pptx_path, source_name)
	else:
		pptx_path, source_name = resolve_input_pptx(input_path, None)
		rows = index_rows(pptx_path, source_name)
	csv_schema.write_slide_csv(output_csv, rows)


#============================================
def index_rows(
	pptx_path: str,
	source_name: str,
) -> list[dict[str, str]]:
	"""
	Index rows from a PPTX path.

	Args:
		pptx_path: Path to PPTX.
		source_name: Source basename for CSV rows.

	Returns:
		list[dict[str, str]]: CSV rows.
	"""
	presentation = pptx.Presentation(pptx_path)
	slide_width = int(getattr(presentation, "slide_width", 0) or 0)
	slide_height = int(getattr(presentation, "slide_height", 0) or 0)
	rows = []
	unsupported_shapes = {}
	layout_errors = {}
	for index, slide in enumerate(presentation.slides, 1):
		title_text = ""
		if slide.shapes.title and slide.shapes.title.text_frame:
			title_text = slide.shapes.title.text_frame.text or ""
		notes_text = extract_notes_text(slide)
		slide_hash, _, _ = pptx_hash.compute_slide_hash_from_slide(
			slide,
			notes_text,
		)
		body_text = extract_body_text(slide)
		asset_types = collect_asset_types(slide)
		layout_type = layout_classifier.classify_layout_type(
			slide,
			slide_width,
			slide_height,
			title_text,
			body_text,
		)
		master_name, layout_warning = resolve_master_name(slide)
		if layout_warning:
			layout_errors[index] = layout_warning
		unsupported = collect_unsupported_shapes(slide)
		if unsupported:
			unsupported_shapes[index] = unsupported
		row = build_slide_row(
			source_name,
			index,
			title_text,
			body_text,
			notes_text,
			slide_hash,
			master_name,
			layout_type,
			asset_types,
		)
		rows.append(row)
	report_index_warnings(unsupported_shapes, layout_errors)
	return rows


#============================================
def main() -> None:
	"""
	Main entry point.
	"""
	args = parse_args()
	output_csv = args.output_csv
	if not output_csv:
		base_name = os.path.splitext(args.input_path)[0]
		output_csv = f"{base_name}.csv"
	index_slides_to_csv(args.input_path, output_csv)


if __name__ == "__main__":
	main()
