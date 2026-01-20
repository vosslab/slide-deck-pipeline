#!/usr/bin/env python3

import argparse
import hashlib
import io
import math
import os
import shutil
import subprocess
import tempfile

# PIP3 modules
import pptx
import pptx.enum.shapes
import pptx.util

# local repo modules
import slide_csv


LAYOUT_HINT_ALIASES = {
	"title_and_object": "title_and_content",
	"title_only": "section_header",
	"two_content": "two_column",
}


#============================================
def parse_args() -> argparse.Namespace:
	"""
	Parse command-line arguments.
	"""
	parser = argparse.ArgumentParser(
		description="Rebuild a PPTX from a merged slide CSV."
	)
	parser.add_argument(
		"-i",
		"--input",
		dest="input_csv",
		required=True,
		help="Input merged CSV path",
	)
	parser.add_argument(
		"-o",
		"--output",
		dest="output_path",
		default="",
		help="Output PPTX or ODP path (default: <input_csv>.pptx)",
	)
	parser.add_argument(
		"-t",
		"--template",
		dest="template_path",
		default="",
		help="Template PPTX path",
	)
	args = parser.parse_args()
	return args


#============================================
def convert_odp_to_pptx(odp_path: str, work_dir: str) -> str:
	"""
	Convert an ODP file to PPTX using soffice.

	Args:
		odp_path: Path to the ODP file.
		work_dir: Output directory for the converted PPTX.

	Returns:
		str: Path to the converted PPTX file.
	"""
	soffice_bin = shutil.which("soffice")
	if not soffice_bin:
		raise FileNotFoundError("soffice not found. Install LibreOffice to convert ODP.")
	command = [
		soffice_bin,
		"--headless",
		"--convert-to",
		"pptx",
		"--outdir",
		work_dir,
		odp_path,
	]
	result = subprocess.run(command, capture_output=True, text=True, cwd=work_dir)
	if result.returncode != 0:
		message = result.stderr.strip() or result.stdout.strip()
		raise RuntimeError(f"ODP conversion failed: {message}")
	base_name = os.path.splitext(os.path.basename(odp_path))[0]
	pptx_path = os.path.join(work_dir, f"{base_name}.pptx")
	if not os.path.exists(pptx_path):
		raise FileNotFoundError(f"Converted PPTX not found: {pptx_path}")
	return pptx_path


#============================================
def convert_pptx_to_odp(pptx_path: str, output_path: str) -> None:
	"""
	Convert a PPTX to ODP using soffice.

	Args:
		pptx_path: Path to PPTX file.
		output_path: Desired ODP output path.
	"""
	soffice_bin = shutil.which("soffice")
	if not soffice_bin:
		raise FileNotFoundError("soffice not found. Install LibreOffice to convert ODP.")
	output_dir = os.path.dirname(output_path) or "."
	command = [
		soffice_bin,
		"--headless",
		"--convert-to",
		"odp",
		"--outdir",
		output_dir,
		pptx_path,
	]
	result = subprocess.run(command, capture_output=True, text=True, cwd=output_dir)
	if result.returncode != 0:
		message = result.stderr.strip() or result.stdout.strip()
		raise RuntimeError(f"PPTX to ODP conversion failed: {message}")
	expected_name = f"{os.path.splitext(os.path.basename(pptx_path))[0]}.odp"
	converted_path = os.path.join(output_dir, expected_name)
	if not os.path.exists(converted_path):
		raise FileNotFoundError(f"Converted ODP not found: {converted_path}")
	if os.path.abspath(converted_path) != os.path.abspath(output_path):
		os.replace(converted_path, output_path)


#============================================
def normalize_layout_name(name: str) -> str:
	"""
	Normalize a layout name to a hint token.

	Args:
		name: Layout name.

	Returns:
		str: Normalized name.
	"""
	if not name:
		return ""
	return name.strip().lower().replace(" ", "_")


#============================================
def select_layout(presentation: pptx.Presentation, layout_hint: str) -> pptx.slide.SlideLayout:
	"""
	Select a slide layout based on a hint.

	Args:
		presentation: Presentation instance.
		layout_hint: Layout hint token.

	Returns:
		pptx.slide.SlideLayout: Selected layout.
	"""
	layout_map = {}
	for layout in presentation.slide_layouts:
		key = normalize_layout_name(layout.name)
		if key and key not in layout_map:
			layout_map[key] = layout
	hint = normalize_layout_name(layout_hint)
	if hint in layout_map:
		return layout_map[hint]
	alias = LAYOUT_HINT_ALIASES.get(hint)
	if alias and alias in layout_map:
		return layout_map[alias]
	if presentation.slide_layouts:
		return presentation.slide_layouts[0]
	raise ValueError("No slide layouts available in template.")


#============================================
def parse_body_lines(body_text: str) -> list[tuple[int, str]]:
	"""
	Parse body text into indentation levels.

	Args:
		body_text: Body text with leading tabs for indentation.

	Returns:
		list[tuple[int, str]]: List of (level, text).
	"""
	if not body_text:
		return []
	lines = []
	cleaned = body_text.replace("\r\n", "\n").replace("\r", "\n")
	for raw_line in cleaned.split("\n"):
		if not raw_line.strip():
			continue
		level = len(raw_line) - len(raw_line.lstrip("\t"))
		text = raw_line.lstrip("\t").strip()
		lines.append((level, text))
	return lines


#============================================
def set_title(slide: pptx.slide.Slide, title_text: str) -> None:
	"""
	Set the slide title if a title placeholder exists.

	Args:
		slide: Slide instance.
		title_text: Title text.
	"""
	if not title_text:
		return
	title_shape = slide.shapes.title
	if not title_shape:
		return
	if not title_shape.text_frame:
		return
	title_shape.text_frame.text = title_text


#============================================
def find_body_placeholder(slide: pptx.slide.Slide) -> pptx.shapes.base.BaseShape | None:
	"""
	Find the primary body placeholder on a slide.

	Args:
		slide: Slide instance.

	Returns:
		pptx.shapes.base.BaseShape | None: Body shape or None.
	"""
	for shape in slide.shapes:
		if not shape.is_placeholder:
			continue
		placeholder_type = shape.placeholder_format.type
		if placeholder_type == pptx.enum.shapes.PP_PLACEHOLDER.BODY:
			return shape
	for shape in slide.shapes:
		if shape.has_text_frame and shape != slide.shapes.title:
			return shape
	return None


#============================================
def set_body_text(slide: pptx.slide.Slide, body_text: str) -> None:
	"""
	Set the body text in the best placeholder.

	Args:
		slide: Slide instance.
		body_text: Body text with indentation markers.
	"""
	lines = parse_body_lines(body_text)
	if not lines:
		return
	body_shape = find_body_placeholder(slide)
	if not body_shape:
		return
	text_frame = body_shape.text_frame
	if not text_frame:
		return
	text_frame.clear()
	for index, (level, text) in enumerate(lines):
		if index == 0:
			paragraph = text_frame.paragraphs[0]
		else:
			paragraph = text_frame.add_paragraph()
		paragraph.text = text
		paragraph.level = level


#============================================
def resolve_source_path(source_pptx: str, csv_dir: str) -> str:
	"""
	Resolve a source path using the CSV directory as fallback.

	Args:
		source_pptx: Source PPTX or ODP path from CSV.
		csv_dir: Directory containing the CSV.

	Returns:
		str: Resolved source path.
	"""
	if os.path.exists(source_pptx):
		return source_pptx
	if csv_dir:
		candidate = os.path.join(csv_dir, source_pptx)
		if os.path.exists(candidate):
			return candidate
	raise FileNotFoundError(f"Source file not found: {source_pptx}")


#============================================
def collect_source_images(slide: pptx.slide.Slide) -> list[dict[str, str | bytes]]:
	"""
	Collect image blobs and hashes from a source slide.

	Args:
		slide: Source slide instance.

	Returns:
		list[dict[str, str | bytes]]: Image records.
	"""
	images = []
	for shape in slide.shapes:
		if shape.shape_type != pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
			continue
		blob = shape.image.blob
		digest = hashlib.sha256(blob).hexdigest()
		shape_id = shape.shape_id
		images.append({"blob": blob, "hash": digest, "shape_id": shape_id})
	return images


#============================================
def sources_match(locator_source: str, row_source: str) -> bool:
	"""
	Check whether locator and row sources refer to the same file.

	Args:
		locator_source: Source from locator.
		row_source: Source from CSV row.

	Returns:
		bool: True if sources match.
	"""
	if locator_source == row_source:
		return True
	if os.path.basename(locator_source) == os.path.basename(row_source):
		return True
	return False


#============================================
def select_images_by_locator(
	images: list[dict[str, str | bytes]],
	image_locators: list[str],
	row_source: str,
	slide_index: int,
) -> list[bytes]:
	"""
	Select images using locator strings.

	Args:
		images: Image records.
		image_locators: Locator strings.
		row_source: Source PPTX or ODP name.
		slide_index: Slide index.

	Returns:
		list[bytes]: Ordered image blobs.
	"""
	if not image_locators:
		return []
	lookup: dict[str, list[bytes]] = {}
	for image in images:
		shape_id = str(image.get("shape_id", ""))
		if not shape_id:
			continue
		lookup.setdefault(shape_id, []).append(image["blob"])
	ordered = []
	for locator in image_locators:
		parsed = slide_csv.parse_image_locator(locator)
		if not parsed:
			continue
		locator_source = parsed.get("source", "")
		if not sources_match(locator_source, row_source):
			continue
		locator_slide = parsed.get("slide", "")
		if locator_slide != str(slide_index):
			continue
		shape_id = parsed.get("shape_id", "")
		if not shape_id:
			continue
		candidates = lookup.get(shape_id, [])
		if not candidates:
			continue
		ordered.append(candidates.pop(0))
	return ordered


#============================================
def select_images_by_hash(
	images: list[dict[str, str | bytes]],
	image_hashes: list[str],
) -> list[bytes]:
	"""
	Select images in the order of hash list.

	Args:
		images: Image records.
		image_hashes: Desired image hashes.

	Returns:
		list[bytes]: Ordered image blobs.
	"""
	if not image_hashes:
		return [image["blob"] for image in images]
	lookup = {}
	for image in images:
		lookup.setdefault(image["hash"], []).append(image["blob"])
	ordered = []
	for digest in image_hashes:
		candidates = lookup.get(digest, [])
		if not candidates:
			continue
		ordered.append(candidates.pop(0))
	return ordered


#============================================
def place_images_grid(slide: pptx.slide.Slide, image_blobs: list[bytes]) -> None:
	"""
	Place images on a slide using a simple grid.

	Args:
		slide: Slide instance.
		image_paths: Image file paths.
	"""
	if not image_blobs:
		return
	cols = 1
	if len(image_blobs) > 1:
		cols = 2
	rows = int(math.ceil(len(image_blobs) / cols))
	margin = pptx.util.Inches(0.5)
	slide_width = slide.part.presentation.slide_width
	slide_height = slide.part.presentation.slide_height
	cell_width = (slide_width - (margin * (cols + 1))) // cols
	cell_height = (slide_height - (margin * (rows + 1))) // rows
	for index, blob in enumerate(image_blobs):
		row = index // cols
		col = index % cols
		left = margin + (cell_width + margin) * col
		top = margin + (cell_height + margin) * row
		stream = io.BytesIO(blob)
		slide.shapes.add_picture(
			stream,
			left,
			top,
			width=cell_width,
			height=cell_height,
		)


#============================================
def insert_images(slide: pptx.slide.Slide, image_blobs: list[bytes]) -> None:
	"""
	Insert images into a slide.

	Args:
		slide: Slide instance.
		image_paths: Image file paths.
	"""
	if not image_blobs:
		return
	picture_placeholders = []
	for shape in slide.shapes:
		if not shape.is_placeholder:
			continue
		if shape.placeholder_format.type == pptx.enum.shapes.PP_PLACEHOLDER.PICTURE:
			picture_placeholders.append(shape)
	if len(image_blobs) == 1 and picture_placeholders:
		stream = io.BytesIO(image_blobs[0])
		picture_placeholders[0].insert_picture(stream)
		return
	place_images_grid(slide, image_blobs)


#============================================
def rebuild_from_csv(
	input_csv: str,
	output_path: str,
	template_path: str,
) -> None:
	"""
	Rebuild a presentation from CSV rows.

	Args:
		input_csv: Input merged CSV path.
		output_path: Output PPTX or ODP path.
		template_path: Template PPTX path or empty string.
	"""
	rows = slide_csv.read_slide_csv(input_csv)
	if template_path:
		presentation = pptx.Presentation(template_path)
	else:
		presentation = pptx.Presentation()
	csv_dir = os.path.dirname(os.path.abspath(input_csv))
	source_cache: dict[str, pptx.Presentation] = {}
	temp_dirs = []
	for row in rows:
		source_pptx = row["source_pptx"]
		source_path = resolve_source_path(source_pptx, csv_dir)
		source_key = source_path
		source_presentation = source_cache.get(source_key)
		if not source_presentation:
			if source_path.lower().endswith(".odp"):
				temp_dir = tempfile.TemporaryDirectory()
				temp_dirs.append(temp_dir)
				converted = convert_odp_to_pptx(source_path, temp_dir.name)
				source_presentation = pptx.Presentation(converted)
			else:
				source_presentation = pptx.Presentation(source_path)
			source_cache[source_key] = source_presentation

		slide_index = int(row["source_slide_index"])
		if slide_index < 1 or slide_index > len(source_presentation.slides):
			raise ValueError(
				f"Source slide index out of range: {source_pptx} {slide_index}."
			)
		source_slide = source_presentation.slides[slide_index - 1]
		source_images = collect_source_images(source_slide)
		image_locators = slide_csv.split_list_field(row["image_locators"])
		image_hashes = slide_csv.split_list_field(row["image_hashes"])
		image_blobs = select_images_by_locator(
			source_images,
			image_locators,
			source_pptx,
			slide_index,
		)
		if not image_blobs:
			image_blobs = select_images_by_hash(source_images, image_hashes)

		layout = select_layout(presentation, row["layout_hint"])
		slide = presentation.slides.add_slide(layout)
		set_title(slide, row["title_text"])
		set_body_text(slide, row["body_text"])
		insert_images(slide, image_blobs)
	if output_path.lower().endswith(".odp"):
		with tempfile.TemporaryDirectory() as temp_dir:
			temp_pptx = os.path.join(temp_dir, "merged.pptx")
			presentation.save(temp_pptx)
			convert_pptx_to_odp(temp_pptx, output_path)
	else:
		presentation.save(output_path)
	for temp_dir in temp_dirs:
		temp_dir.cleanup()


#============================================
def main() -> None:
	"""
	Main entry point.
	"""
	args = parse_args()
	output_path = args.output_path
	if not output_path:
		base_name = os.path.splitext(args.input_csv)[0]
		output_path = f"{base_name}.pptx"
	rebuild_from_csv(
		args.input_csv,
		output_path,
		args.template_path,
	)


if __name__ == "__main__":
	main()
