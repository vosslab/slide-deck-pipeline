# Standard Library
import os
import tempfile

# PIP3 modules
import pptx

# local repo modules
import slide_deck_pipeline.pptx_io as pptx_io
import slide_deck_pipeline.soffice_tools as soffice_tools


#============================================
def clear_direct_formatting(
	input_path: str,
	output_path: str,
	inplace: bool,
) -> tuple[int, int]:
	"""
	Clear direct formatting from all text boxes in a PPTX or ODP file.

	Args:
		input_path: Input PPTX or ODP path.
		output_path: Output PPTX or ODP path.
		inplace: Allow writing to the input path.

	Returns:
		tuple[int, int]: (text runs seen, text runs cleared).
	"""
	if not inplace:
		input_abs = os.path.abspath(input_path)
		output_abs = os.path.abspath(output_path)
		if input_abs == output_abs:
			raise ValueError("Output path matches input; use --inplace to override.")
	output_is_odp = output_path.lower().endswith(".odp")
	if input_path.lower().endswith(".odp"):
		with tempfile.TemporaryDirectory() as temp_dir:
			pptx_path, _ = pptx_io.resolve_input_pptx(input_path, temp_dir)
			return clear_pptx(pptx_path, output_path, output_is_odp)
	pptx_path, _ = pptx_io.resolve_input_pptx(input_path, None)
	return clear_pptx(pptx_path, output_path, output_is_odp)


#============================================
def clear_pptx(
	pptx_path: str,
	output_path: str,
	output_is_odp: bool,
) -> tuple[int, int]:
	"""
	Clear direct formatting from all text boxes in a PPTX file.

	Args:
		pptx_path: Input PPTX path.
		output_path: Output path.
		output_is_odp: True if output should be ODP.

	Returns:
		tuple[int, int]: (text runs seen, text runs cleared).
	"""
	presentation = pptx.Presentation(pptx_path)
	total = 0
	cleared = 0
	for slide in presentation.slides:
		for shape in slide.shapes:
			if not shape.has_text_frame:
				continue
			text_frame = shape.text_frame
			for paragraph in text_frame.paragraphs:
				for run in paragraph.runs:
					total += 1
					# Check if run has direct formatting
					had_formatting = False

					# Clear font properties
					if run.font.name is not None:
						run.font.name = None
						had_formatting = True
					if run.font.size is not None:
						run.font.size = None
						had_formatting = True
					if run.font.bold is not None:
						run.font.bold = None
						had_formatting = True
					if run.font.italic is not None:
						run.font.italic = None
						had_formatting = True
					if run.font.underline is not None:
						run.font.underline = None
						had_formatting = True

					# Clear font color if it's set
					try:
						if run.font.color.type is not None:
							# Setting to None doesn't always work, so we skip color for now
							# This is a limitation of python-pptx
							pass
					except Exception:
						pass

					if had_formatting:
						cleared += 1

	if output_is_odp:
		with tempfile.TemporaryDirectory() as temp_dir:
			temp_pptx = os.path.join(temp_dir, "format_cleared.pptx")
			presentation.save(temp_pptx)
			soffice_tools.convert_pptx_to_odp(temp_pptx, output_path)
	else:
		presentation.save(output_path)
	return (total, cleared)
