# PIP3 modules
import pptx


#============================================
def reset_picture_crop(shape) -> None:
	"""
	Reset crop settings on a picture shape when supported.

	Args:
		shape: Picture shape.
	"""
	for attr_name in ("crop_left", "crop_right", "crop_top", "crop_bottom"):
		if hasattr(shape, attr_name):
			setattr(shape, attr_name, 0)


#============================================
def fit_picture_shape(
	shape,
	box_left: int,
	box_top: int,
	box_width: int,
	box_height: int,
) -> bool:
	"""
	Fit a picture shape inside a bounding box without stretching.

	Args:
		shape: Picture shape.
		box_left: Bounding box left.
		box_top: Bounding box top.
		box_width: Bounding box width.
		box_height: Bounding box height.

	Returns:
		bool: True if the shape was adjusted.
	"""
	if box_width <= 0 or box_height <= 0:
		return False
	if not hasattr(shape, "image"):
		return False
	image_size = getattr(shape.image, "size", None)
	if not image_size:
		return False
	image_width, image_height = image_size
	if not image_width or not image_height:
		return False
	image_ratio = image_width / image_height
	box_ratio = box_width / box_height
	if box_ratio >= image_ratio:
		new_height = box_height
		new_width = int(round(box_height * image_ratio))
	else:
		new_width = box_width
		new_height = int(round(box_width / image_ratio))
	if new_width <= 0 or new_height <= 0:
		return False
	if new_width > box_width:
		new_width = box_width
	if new_height > box_height:
		new_height = box_height
	left = int(box_left + (box_width - new_width) / 2)
	top = int(box_top + (box_height - new_height) / 2)
	reset_picture_crop(shape)
	shape.left = left
	shape.top = top
	shape.width = new_width
	shape.height = new_height
	return True


#============================================
def iter_picture_shapes(slide: pptx.slide.Slide):
	"""
	Yield picture shapes from a slide, including grouped shapes.

	Args:
		slide: Slide instance.

	Yields:
		Picture shapes.
	"""
	for shape in iter_shapes(slide.shapes):
		if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
			yield shape


#============================================
def iter_shapes(shapes):
	"""
	Yield shapes recursively from a shape collection.

	Args:
		shapes: Shape collection.

	Yields:
		Shape objects.
	"""
	for shape in shapes:
		if (
			getattr(shape, "shape_type", None)
			== pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP
			and hasattr(shape, "shapes")
		):
			for nested in iter_shapes(shape.shapes):
				yield nested
			continue
		yield shape
