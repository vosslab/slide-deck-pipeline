# PIP3 modules
import pptx
import pptx.enum.shapes


#============================================
def classify_placeholder_role(placeholder_type) -> str | None:
	"""
	Map a placeholder type to a semantic role.

	Args:
		placeholder_type: pptx placeholder type enum value.

	Returns:
		str | None: Role name or None if unsupported.
	"""
	placeholders = pptx.enum.shapes.PP_PLACEHOLDER
	if placeholder_type in (
		placeholders.TITLE,
		placeholders.CENTER_TITLE,
	):
		return "title"
	if placeholder_type == placeholders.SUBTITLE:
		return "subtitle"
	body_types = [
		placeholders.BODY,
		placeholders.OBJECT,
	]
	content_type = getattr(placeholders, "CONTENT", None)
	if content_type is not None:
		body_types.append(content_type)
	text_type = getattr(placeholders, "TEXT", None)
	if text_type is not None:
		body_types.append(text_type)
	if placeholder_type in tuple(body_types):
		return "body"
	return None


#============================================
def collect_placeholder_boxes(slide: pptx.slide.Slide) -> list[dict[str, object]]:
	"""
	Collect placeholder boxes with roles and geometry.

	Args:
		slide: Slide instance.

	Returns:
		list[dict[str, object]]: Placeholder metadata.
	"""
	boxes = []
	for shape in slide.shapes:
		if not getattr(shape, "is_placeholder", False):
			continue
		try:
			placeholder_type = shape.placeholder_format.type
		except Exception:
			continue
		role = classify_placeholder_role(placeholder_type)
		if not role:
			continue
		boxes.append(
			{
				"role": role,
				"left": int(shape.left),
				"top": int(shape.top),
				"width": int(shape.width),
				"height": int(shape.height),
			}
		)
	return boxes


#============================================
def is_two_content_split(
	body_boxes: list[dict[str, object]],
	slide_width: int,
) -> bool:
	"""
	Check if two body placeholders are split left-right.

	Args:
		body_boxes: Body placeholder metadata.
		slide_width: Slide width in EMUs.

	Returns:
		bool: True if split left-right.
	"""
	if len(body_boxes) != 2 or slide_width <= 0:
		return False
	sorted_boxes = sorted(body_boxes, key=lambda box: box["left"])
	first = sorted_boxes[0]
	second = sorted_boxes[1]
	left = int(first["left"])
	right = int(second["left"])
	width_left = int(first["width"])
	width_right = int(second["width"])
	top_left = int(first["top"])
	top_right = int(second["top"])
	height_left = int(first["height"])
	height_right = int(second["height"])
	center_diff_y = abs((top_left + height_left / 2) - (top_right + height_right / 2))
	if max(height_left, height_right) == 0:
		return False
	if center_diff_y > max(height_left, height_right) * 0.2:
		return False
	width_ratio = min(width_left, width_right) / max(width_left, width_right)
	if width_ratio < 0.6:
		return False
	horizontal_overlap = min(left + width_left, right + width_right) - max(left, right)
	if horizontal_overlap > min(width_left, width_right) * 0.2:
		return False
	if left > right:
		return False
	return True


#============================================
def is_centered_box(
	box: dict[str, object],
	slide_width: int,
	slide_height: int,
) -> bool:
	"""
	Check if a box is centered on the slide.

	Args:
		box: Box metadata.
		slide_width: Slide width in EMUs.
		slide_height: Slide height in EMUs.

	Returns:
		bool: True if box is centered.
	"""
	if slide_width <= 0 or slide_height <= 0:
		return False
	left = int(box["left"])
	top = int(box["top"])
	width = int(box["width"])
	height = int(box["height"])
	center_x = left + width / 2
	center_y = top + height / 2
	if abs(center_x - slide_width / 2) > slide_width * 0.1:
		return False
	if abs(center_y - slide_height / 2) > slide_height * 0.15:
		return False
	if width > slide_width * 0.9:
		return False
	return True


#============================================
def classify_layout_type(
	slide: pptx.slide.Slide,
	slide_width: int,
	slide_height: int,
	title_text: str,
	body_text: str,
) -> tuple[str, float, list[str]]:
	"""
	Classify a slide into a semantic layout type.

	Args:
		slide: Slide instance.
		slide_width: Slide width in EMUs.
		slide_height: Slide height in EMUs.
		title_text: Title text.
		body_text: Body text.

	Returns:
		tuple[str, float, list[str]]: Layout type, confidence, reasons.
	"""
	placeholders = collect_placeholder_boxes(slide)
	if not placeholders:
		if title_text or body_text:
			return ("custom", 0.2, ["text_without_placeholders"])
		return ("blank", 1.0, ["no_placeholders_no_text"])
	title_boxes = [box for box in placeholders if box["role"] == "title"]
	subtitle_boxes = [box for box in placeholders if box["role"] == "subtitle"]
	body_boxes = [box for box in placeholders if box["role"] == "body"]
	if title_boxes and subtitle_boxes and not body_boxes:
		return ("title_slide", 1.0, ["title_and_subtitle"])
	if title_boxes and not subtitle_boxes and not body_boxes:
		return ("title_only", 0.9, ["title_only"])
	if title_boxes and len(body_boxes) == 1:
		return ("title_content", 0.9, ["title_and_body"])
	if title_boxes and len(body_boxes) == 2:
		if is_two_content_split(body_boxes, slide_width):
			return ("two_content", 0.9, ["two_body_split"])
		return ("custom", 0.4, ["two_body_not_split"])
	if not title_boxes and len(body_boxes) == 1:
		if is_centered_box(body_boxes[0], slide_width, slide_height):
			return ("centered_text", 0.8, ["centered_body"])
	return ("custom", 0.2, ["unclassified"])
