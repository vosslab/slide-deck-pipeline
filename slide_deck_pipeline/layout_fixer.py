# Standard Library
import os
import re
import tempfile

# PIP3 modules
import pptx
import pptx.enum.shapes

# local repo modules
import slide_deck_pipeline.pptx_io as pptx_io
import slide_deck_pipeline.text_boxes as text_boxes
import slide_deck_pipeline.soffice_tools as soffice_tools


# Title length limit in characters
TITLE_MAX_LENGTH = 120


#============================================
def fix_layout(
	input_path: str,
	output_path: str,
	inplace: bool,
	verbose: bool = False,
) -> tuple[int, int, int]:
	"""
	Fix slide layouts by rearranging text content based on sensible rules.

	Args:
		input_path: Input PPTX or ODP path.
		output_path: Output PPTX or ODP path.
		inplace: Allow writing to the input path.

	Returns:
		tuple[int, int, int]: (slides inspected, swaps made, moves made).
	"""
	if not inplace:
		input_abs = os.path.abspath(input_path)
		output_abs = os.path.abspath(output_path)
		if input_abs == output_abs:
			raise ValueError("Output path matches input; use --inplace to override.")
	output_is_odp = output_path.lower().endswith(".odp")
	if input_path.lower().endswith(".odp"):
		with tempfile.TemporaryDirectory() as temp_dir:
			pptx_path, _ = pptx_io.resolve_input_pptx(input_path, temp_dir)
			return fix_pptx(pptx_path, output_path, output_is_odp, verbose)
	pptx_path, _ = pptx_io.resolve_input_pptx(input_path, None)
	return fix_pptx(pptx_path, output_path, output_is_odp, verbose)


#============================================
def is_title_placeholder(shape) -> bool:
	"""Check if shape is a title placeholder."""
	if not shape.is_placeholder:
		return False
	try:
		return shape.placeholder_format.type == pptx.enum.shapes.PP_PLACEHOLDER.TITLE
	except Exception:
		return False


#============================================
def is_body_placeholder(shape) -> bool:
	"""Check if shape is a body/content placeholder."""
	if not shape.is_placeholder:
		return False
	try:
		placeholder_type = shape.placeholder_format.type
		placeholders = pptx.enum.shapes.PP_PLACEHOLDER
		body_types = [placeholders.BODY]
		for attr_name in ("OBJECT", "CONTENT", "TEXT"):
			candidate = getattr(placeholders, attr_name, None)
			if candidate is not None:
				body_types.append(candidate)
		return placeholder_type in tuple(body_types)
	except Exception:
		return False


#============================================
def get_text_length(shape) -> int:
	"""Get the character count of text in a shape."""
	if not shape.has_text_frame:
		return 0
	text = text_boxes.extract_text_block(shape)
	return len(text.strip())


#============================================
def set_text(shape, text: str) -> None:
	"""Set the text content of a shape, preserving structure."""
	if not shape.has_text_frame:
		return
	text_frame = shape.text_frame
	text_frame.clear()
	text_frame.text = text


#============================================
def get_text(shape) -> str:
	"""Get the text content of a shape."""
	if not shape.has_text_frame:
		return ""
	return text_boxes.extract_text_block(shape).strip()


#============================================
def make_short_title(long_text: str) -> str:
	"""Create a short title from long text."""
	# Take first sentence or first line
	lines = long_text.strip().split('\n')
	first_line = lines[0].strip()

	# Try to find a sentence break
	sentences = re.split(r'[.!?]', first_line)
	if sentences and len(sentences[0].strip()) <= TITLE_MAX_LENGTH:
		return sentences[0].strip()

	# Take first TITLE_MAX_LENGTH chars and add ellipsis
	if len(first_line) > TITLE_MAX_LENGTH:
		return first_line[:TITLE_MAX_LENGTH - 3].strip() + "..."

	return first_line


#============================================
def count_assets(slide) -> dict[str, int]:
	"""Count different types of assets on a slide."""
	counts = {
		"images": 0,
		"tables": 0,
		"charts": 0,
		"other": 0,
	}
	for shape in slide.shapes:
		if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
			counts["images"] += 1
		elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
			counts["tables"] += 1
		elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.CHART:
			counts["charts"] += 1
		elif not shape.has_text_frame and shape.shape_type != pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
			counts["other"] += 1
	return counts


#============================================
def fix_slide_layout(slide, slide_num: int) -> tuple[bool, bool, str]:
	"""
	Fix layout issues on a single slide.

	Returns:
		tuple[bool, bool, str]: (swapped, moved, description)
	"""
	# Find title and body placeholders
	title_shape = None
	body_shapes = []

	for shape in slide.shapes:
		if is_title_placeholder(shape):
			title_shape = shape
		elif is_body_placeholder(shape):
			body_shapes.append(shape)

	if not title_shape:
		return (False, False, "No title placeholder")

	title_text = get_text(title_shape)
	title_len = len(title_text)

	# Count assets
	assets = count_assets(slide)
	asset_desc = []
	if assets["images"] > 0:
		asset_desc.append(f"{assets['images']} image(s)")
	if assets["tables"] > 0:
		asset_desc.append(f"{assets['tables']} table(s)")
	if assets["charts"] > 0:
		asset_desc.append(f"{assets['charts']} chart(s)")

	# No title text - nothing to fix
	if title_len == 0:
		return (False, False, f"Empty title, assets: {', '.join(asset_desc) if asset_desc else 'none'}")

	# Get primary body shape and text
	body_shape = body_shapes[0] if body_shapes else None
	body_text = get_text(body_shape) if body_shape else ""
	body_len = len(body_text)

	# Rule 1: If title is longer than body (and both have text), swap them
	if body_len > 0 and title_len > body_len * 1.5 and title_len > TITLE_MAX_LENGTH:
		set_text(title_shape, body_text)
		set_text(body_shape, title_text)
		return (True, False, f"Swapped title ({title_len} chars) with body ({body_len} chars)")

	# Rule 2: If title is too long and body exists, move title to body and create short title
	if title_len > TITLE_MAX_LENGTH and body_shape:
		if body_len == 0:
			# Body is empty - move title to body, create short title from it
			short_title = make_short_title(title_text)
			set_text(body_shape, title_text)
			set_text(title_shape, short_title)
			return (False, True, f"Moved long title ({title_len} chars) to body, created short title ({len(short_title)} chars)")
		elif body_len < title_len:
			# Body has some text but less than title - swap them
			set_text(title_shape, body_text)
			set_text(body_shape, title_text)
			return (True, False, f"Swapped long title ({title_len} chars) with shorter body ({body_len} chars)")

	# Rule 3: Title is too long but no body to move it to - just report
	if title_len > TITLE_MAX_LENGTH and not body_shape:
		return (False, False, f"Title too long ({title_len} chars) but no body placeholder available")

	return (False, False, f"OK: title={title_len} chars, body={body_len} chars, assets: {', '.join(asset_desc) if asset_desc else 'none'}")


#============================================
def fix_pptx(
	pptx_path: str,
	output_path: str,
	output_is_odp: bool,
	verbose: bool = False,
) -> tuple[int, int, int]:
	"""
	Fix layouts in a PPTX file.

	Args:
		pptx_path: Input PPTX path.
		output_path: Output path.
		output_is_odp: True if output should be ODP.

	Returns:
		tuple[int, int, int]: (slides inspected, swaps made, moves made).
	"""
	presentation = pptx.Presentation(pptx_path)
	total_slides = 0
	total_swaps = 0
	total_moves = 0

	for slide_num, slide in enumerate(presentation.slides, start=1):
		total_slides += 1
		swapped, moved, description = fix_slide_layout(slide, slide_num)

		if swapped:
			total_swaps += 1
			print(f"  Slide {slide_num}: {description}")
		elif moved:
			total_moves += 1
			print(f"  Slide {slide_num}: {description}")
		elif verbose:
			print(f"  Slide {slide_num}: {description}")

	if output_is_odp:
		with tempfile.TemporaryDirectory() as temp_dir:
			temp_pptx = os.path.join(temp_dir, "layout_fixed.pptx")
			presentation.save(temp_pptx)
			soffice_tools.convert_pptx_to_odp(temp_pptx, output_path)
	else:
		presentation.save(output_path)
	return (total_slides, total_swaps, total_moves)
