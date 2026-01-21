# Standard Library
import hashlib

# PIP3 modules
import pptx
import pptx.enum.shapes

# local repo modules
import slide_deck_pipeline.csv_schema as csv_schema
import slide_deck_pipeline.layout_classifier as layout_classifier
import slide_deck_pipeline.pptx_text as pptx_text


#============================================
def extract_slide_xml(slide) -> bytes:
	"""
	Return slide XML bytes.

	Args:
		slide: Slide instance.

	Returns:
		bytes: Slide XML.
	"""
	return slide.part.blob


#============================================
def shape_type_name(shape) -> str:
	"""
	Return a readable shape type name.

	Args:
		shape: Shape instance.

	Returns:
		str: Shape type name.
	"""
	shape_type = getattr(shape, "shape_type", None)
	if shape_type is None:
		return "unknown"
	return getattr(shape_type, "name", str(shape_type))


#============================================
def shape_geometry(shape) -> tuple[int, int, int, int]:
	"""
	Return shape geometry as integer tuple.

	Args:
		shape: Shape instance.

	Returns:
		tuple[int, int, int, int]: (left, top, width, height).
	"""
	return (
		int(getattr(shape, "left", 0) or 0),
		int(getattr(shape, "top", 0) or 0),
		int(getattr(shape, "width", 0) or 0),
		int(getattr(shape, "height", 0) or 0),
	)


#============================================
def placeholder_role(shape) -> str:
	"""
	Return placeholder role for a shape.

	Args:
		shape: Shape instance.

	Returns:
		str: Placeholder role name or empty string.
	"""
	if not getattr(shape, "is_placeholder", False):
		return ""
	try:
		placeholder_type = shape.placeholder_format.type
	except Exception:
		return ""
	return layout_classifier.classify_placeholder_role(placeholder_type) or ""


#============================================
def hash_bytes(payload: bytes) -> str:
	"""
	Hash bytes to a stable short digest.

	Args:
		payload: Input bytes.

	Returns:
		str: Short hash string.
	"""
	return hashlib.sha256(payload).hexdigest()[:16]


#============================================
def hash_image_blob(shape) -> str:
	"""
	Hash a picture blob if present.

	Args:
		shape: Shape instance.

	Returns:
		str: Image hash or empty string.
	"""
	if getattr(shape, "shape_type", None) != pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
		return ""
	image = getattr(shape, "image", None)
	if not image or not getattr(image, "blob", None):
		return ""
	return hash_bytes(bytes(image.blob))


#============================================
def hash_shape_text(shape) -> str:
	"""
	Hash normalized shape text.

	Args:
		shape: Shape instance.

	Returns:
		str: Text hash or empty string.
	"""
	lines = pptx_text.extract_shape_text(shape)
	if not lines:
		return ""
	text = "\n".join(lines)
	return csv_schema.compute_text_hash(text)


#============================================
def shape_kind(shape) -> str:
	"""
	Return a coarse shape kind.

	Args:
		shape: Shape instance.

	Returns:
		str: Shape kind.
	"""
	if (
		getattr(shape, "shape_type", None)
		== pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP
		and hasattr(shape, "shapes")
	):
		return "group"
	if getattr(shape, "shape_type", None) == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
		return "picture"
	if getattr(shape, "has_table", False):
		return "table"
	if getattr(shape, "has_chart", False):
		return "chart"
	if getattr(shape, "is_placeholder", False):
		return "placeholder"
	if getattr(shape, "has_text_frame", False):
		return "textbox"
	return "shape"


#============================================
def build_shape_tokens(shape, tokens: list[tuple]) -> None:
	"""
	Append shape tokens in order.

	Args:
		shape: Shape instance.
		tokens: Token list to append to.
	"""
	kind = shape_kind(shape)
	geom = shape_geometry(shape)
	if kind == "group":
		tokens.append(("group_start", geom, shape_type_name(shape)))
		for nested in shape.shapes:
			build_shape_tokens(nested, tokens)
		tokens.append(("group_end",))
		return
	role = placeholder_role(shape)
	text_hash = hash_shape_text(shape)
	image_hash = hash_image_blob(shape)
	tokens.append(
		(
			"shape",
			kind,
			role,
			geom,
			text_hash,
			image_hash,
			shape_type_name(shape),
		)
	)


#============================================
def compute_slide_hash_from_slide(
	slide,
	notes_text: str | None = None,
) -> tuple[str, str, bytes]:
	"""
	Compute slide hash and return slide XML and notes text.

	Args:
		slide: Slide instance.
		notes_text: Optional notes text to reuse.

	Returns:
		tuple[str, str, bytes]: Slide hash, notes text, slide XML bytes.
	"""
	if notes_text is None:
		notes_text = pptx_text.extract_notes_text(slide)
	slide_xml = extract_slide_xml(slide)
	tokens: list[tuple] = []
	for shape in slide.shapes:
		build_shape_tokens(shape, tokens)
	payload = repr(tuple(tokens)).encode("utf-8")
	slide_hash = csv_schema.compute_slide_hash(payload, notes_text)
	return (slide_hash, notes_text, slide_xml)
