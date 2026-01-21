# PIP3 modules
import pptx
import pptx.enum.shapes


#============================================
def extract_paragraph_lines(text_frame: pptx.text.text.TextFrame) -> list[str]:
	"""
	Extract text frame paragraphs with indentation markers.

	Args:
		text_frame: TextFrame instance.

	Returns:
		list[str]: Lines with leading tab indentation.
	"""
	lines = []
	for paragraph in text_frame.paragraphs:
		text = paragraph.text.strip()
		if not text:
			continue
		indent = "\t" * paragraph.level
		lines.append(f"{indent}{text}")
	return lines


#============================================
def extract_table_text(table: pptx.table.Table) -> list[str]:
	"""
	Extract text from a table.

	Args:
		table: Table instance.

	Returns:
		list[str]: Table cell lines.
	"""
	lines = []
	for row in table.rows:
		for cell in row.cells:
			if not cell.text_frame:
				continue
			lines.extend(extract_paragraph_lines(cell.text_frame))
	return lines


#============================================
def extract_chart_title_text(chart) -> list[str]:
	"""
	Extract chart title text when available.

	Args:
		chart: Chart instance.

	Returns:
		list[str]: Chart title lines.
	"""
	if not chart:
		return []
	if not getattr(chart, "has_title", False):
		return []
	chart_title = getattr(chart, "chart_title", None)
	if not chart_title or not getattr(chart_title, "text_frame", None):
		return []
	return extract_paragraph_lines(chart_title.text_frame)


#============================================
def extract_shape_text(shape) -> list[str]:
	"""
	Extract text from a shape, including tables and chart titles.

	Args:
		shape: Shape instance.

	Returns:
		list[str]: Text lines.
	"""
	lines = []
	if (
		getattr(shape, "shape_type", None)
		== pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP
		and hasattr(shape, "shapes")
	):
		for nested in shape.shapes:
			lines.extend(extract_shape_text(nested))
		return lines
	if getattr(shape, "has_text_frame", False):
		lines.extend(extract_paragraph_lines(shape.text_frame))
	if getattr(shape, "has_table", False):
		lines.extend(extract_table_text(shape.table))
	if getattr(shape, "has_chart", False):
		lines.extend(extract_chart_title_text(shape.chart))
	return lines


#============================================
def extract_body_text(slide: pptx.slide.Slide) -> str:
	"""
	Extract body text from non-title text frames.

	Args:
		slide: Slide instance.

	Returns:
		str: Body text with newline separators.
	"""
	lines = []
	title_shape = slide.shapes.title
	for shape in slide.shapes:
		if title_shape and shape == title_shape:
			continue
		lines.extend(extract_shape_text(shape))
	return "\n".join(lines)


#============================================
def extract_notes_text(slide: pptx.slide.Slide) -> str:
	"""
	Extract speaker notes text from a slide.

	Args:
		slide: Slide instance.

	Returns:
		str: Notes text.
	"""
	if not slide.has_notes_slide:
		return ""
	notes_frame = slide.notes_slide.notes_text_frame
	if not notes_frame:
		return ""
	return notes_frame.text or ""


#============================================
def extract_slide_text(slide: pptx.slide.Slide) -> str:
	"""
	Extract all on-slide text for hashing.

	Args:
		slide: Slide instance.

	Returns:
		str: Slide text.
	"""
	lines = []
	for shape in slide.shapes:
		lines.extend(extract_shape_text(shape))
	return "\n".join(lines)
