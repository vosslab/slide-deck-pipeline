# Standard Library
import io
import math
import os
import tempfile

# PIP3 modules
import pptx
import pptx.enum.shapes
import pptx.util

# local repo modules
import slide_deck_pipeline.csv_schema as csv_schema
import slide_deck_pipeline.layout_classifier as layout_classifier
import slide_deck_pipeline.path_resolver as path_resolver
import slide_deck_pipeline.pptx_hash as pptx_hash
import slide_deck_pipeline.soffice_tools as soffice_tools
import slide_deck_pipeline.text_normalization as text_normalization
import slide_deck_pipeline.image_utils as image_utils


#============================================
def normalize_name(name: str) -> str:
	"""
	Normalize a master or layout name for matching.

	Args:
		name: Name value.

	Returns:
		str: Normalized name.
	"""
	return text_normalization.normalize_simple_name(name)


#============================================
def build_layout_map(
	presentation: pptx.Presentation,
) -> dict[str, dict[str, pptx.slide.SlideLayout]]:
	"""
	Build a map of master -> layout_type -> slide layout.

	Args:
		presentation: Presentation instance.

	Returns:
		dict[str, dict[str, pptx.slide.SlideLayout]]: Layout map.
	"""
	slide_width = int(getattr(presentation, "slide_width", 0) or 0)
	slide_height = int(getattr(presentation, "slide_height", 0) or 0)
	layout_map: dict[str, dict[str, pptx.slide.SlideLayout]] = {}
	for layout in presentation.slide_layouts:
		layout_type, _, _ = layout_classifier.classify_layout_type(
			layout,
			slide_width,
			slide_height,
			"",
			"",
		)
		master = getattr(layout, "slide_master", None)
		master_key = normalize_name(getattr(master, "name", "")) or "custom"
		layout_map.setdefault(master_key, {})
		if layout_type not in layout_map[master_key]:
			layout_map[master_key][layout_type] = layout
	return layout_map


#============================================
def select_layout(
	presentation: pptx.Presentation,
	layout_map: dict[str, dict[str, pptx.slide.SlideLayout]],
	master_name: str,
	layout_type: str,
) -> pptx.slide.SlideLayout:
	"""
	Select a slide layout based on master and layout type.

	Args:
		presentation: Presentation instance.
		layout_map: Master/layout_type map.
		master_name: Template master name.
		layout_type: Semantic layout type.

	Returns:
		pptx.slide.SlideLayout: Selected layout.
	"""
	target_master = normalize_name(master_name) or "custom"
	target_layout = normalize_name(layout_type) or "custom"
	master_layouts = layout_map.get(target_master, {})
	if target_layout in master_layouts:
		return master_layouts[target_layout]
	if "custom" in master_layouts:
		return master_layouts["custom"]
	if master_layouts:
		return list(master_layouts.values())[0]
	if presentation.slide_layouts:
		return presentation.slide_layouts[0]
	raise ValueError("No slide layouts available in template.")


#============================================
def parse_body_lines(body_text: str) -> list[tuple[int, str]]:
	"""
	Parse body text into indentation levels.

	Args:
		body_text: Body text with leading tabs for indentation.

	Returns:
		list[tuple[int, str]]: List of (level, text).
	"""
	return text_normalization.parse_tab_indented_lines(
		body_text,
		keep_blank_lines=False,
		strip_text=True,
	)


#============================================
def set_title(slide: pptx.slide.Slide, title_text: str) -> None:
	"""
	Set the slide title if a title placeholder exists.

	Args:
		slide: Slide instance.
		title_text: Title text.
	"""
	if not title_text:
		return
	title_shape = slide.shapes.title
	if not title_shape:
		return
	if not title_shape.text_frame:
		return
	title_shape.text_frame.text = title_text


#============================================
def find_body_placeholder(slide: pptx.slide.Slide) -> pptx.shapes.base.BaseShape | None:
	"""
	Find the primary body placeholder on a slide.

	Args:
		slide: Slide instance.

	Returns:
		pptx.shapes.base.BaseShape | None: Body shape or None.
	"""
	for shape in slide.shapes:
		if not shape.is_placeholder:
			continue
		placeholder_type = shape.placeholder_format.type
		if placeholder_type == pptx.enum.shapes.PP_PLACEHOLDER.BODY:
			return shape
	for shape in slide.shapes:
		if shape.has_text_frame and shape != slide.shapes.title:
			return shape
	return None


#============================================
def set_body_text(slide: pptx.slide.Slide, body_text: str) -> None:
	"""
	Set the body text in the best placeholder.

	Args:
		slide: Slide instance.
		body_text: Body text with indentation markers.
	"""
	lines = parse_body_lines(body_text)
	if not lines:
		return
	body_shape = find_body_placeholder(slide)
	if not body_shape:
		return
	text_frame = body_shape.text_frame
	if not text_frame:
		return
	text_frame.clear()
	for index, (level, text) in enumerate(lines):
		if index == 0:
			paragraph = text_frame.paragraphs[0]
		else:
			paragraph = text_frame.add_paragraph()
		paragraph.text = text
		paragraph.level = level


#============================================
def set_notes_text(slide: pptx.slide.Slide, notes_text: str) -> None:
	"""
	Set speaker notes text on a slide.

	Args:
		slide: Slide instance.
		notes_text: Notes text.
	"""
	if not notes_text:
		return
	notes_slide = slide.notes_slide
	if not notes_slide:
		return
	notes_frame = notes_slide.notes_text_frame
	if not notes_frame:
		return
	notes_frame.text = notes_text


#============================================
def collect_source_images(slide: pptx.slide.Slide) -> list[bytes]:
	"""
	Collect image blobs from a source slide.

	Args:
		slide: Source slide instance.

	Returns:
		list[bytes]: Image blobs in slide order.
	"""
	images = []
	for shape in slide.shapes:
		if shape.shape_type != pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
			continue
		images.append(shape.image.blob)
	return images


#============================================
def place_images_grid(slide: pptx.slide.Slide, image_blobs: list[bytes]) -> None:
	"""
	Place images on a slide using a simple grid.

	Args:
		slide: Slide instance.
		image_blobs: Image blobs.
	"""
	if not image_blobs:
		return
	cols = 1
	if len(image_blobs) > 1:
		cols = 2
	rows = int(math.ceil(len(image_blobs) / cols))
	margin = pptx.util.Inches(0.5)
	slide_width, slide_height = get_slide_dimensions(slide)
	cell_width = (slide_width - (margin * (cols + 1))) // cols
	cell_height = (slide_height - (margin * (rows + 1))) // rows
	for index, blob in enumerate(image_blobs):
		row = index // cols
		col = index % cols
		left = margin + (cell_width + margin) * col
		top = margin + (cell_height + margin) * row
		stream = io.BytesIO(blob)
		picture = slide.shapes.add_picture(
			stream,
			left,
			top,
			width=cell_width,
			height=cell_height,
		)
		image_utils.fit_picture_shape(
			picture,
			int(left),
			int(top),
			int(cell_width),
			int(cell_height),
		)


#============================================
def insert_images(slide: pptx.slide.Slide, image_blobs: list[bytes]) -> None:
	"""
	Insert images into a slide.

	Args:
		slide: Slide instance.
		image_blobs: Image blobs.
	"""
	if not image_blobs:
		return
	picture_placeholders = []
	placeholder_enum = pptx.enum.shapes.PP_PLACEHOLDER
	picture_types = [placeholder_enum.PICTURE]
	for attr_name in ("OBJECT", "CONTENT"):
		candidate = getattr(placeholder_enum, attr_name, None)
		if candidate is not None:
			picture_types.append(candidate)
	for shape in slide.shapes:
		if not shape.is_placeholder:
			continue
		if shape.placeholder_format.type in tuple(picture_types):
			picture_placeholders.append(shape)
	if len(image_blobs) == 1 and picture_placeholders:
		stream = io.BytesIO(image_blobs[0])
		placeholder = picture_placeholders[0]
		picture = placeholder.insert_picture(stream)
		if all(
			hasattr(placeholder, attr)
			for attr in ("left", "top", "width", "height")
		):
			image_utils.fit_picture_shape(
				picture,
				int(placeholder.left),
				int(placeholder.top),
				int(placeholder.width),
				int(placeholder.height),
			)
		return
	place_images_grid(slide, image_blobs)


#============================================
def get_slide_dimensions(slide: pptx.slide.Slide) -> tuple[int, int]:
	"""
	Return slide width and height.

	Args:
		slide: Slide instance.

	Returns:
		tuple[int, int]: Slide width and height.
	"""
	presentation = None
	if hasattr(slide.part, "presentation"):
		presentation = slide.part.presentation
	elif hasattr(slide.part, "package"):
		presentation_part = slide.part.package.presentation_part
		presentation = presentation_part.presentation
	if not presentation:
		raise ValueError("Unable to resolve presentation dimensions.")
	return (presentation.slide_width, presentation.slide_height)


#============================================
def rebuild_from_csv(
	input_csv: str,
	output_path: str,
	template_path: str,
) -> None:
	"""
	Rebuild a presentation from CSV rows.

	Args:
		input_csv: Input merged CSV path.
		output_path: Output PPTX or ODP path.
		template_path: Template PPTX path or empty string.
	"""
	rows = csv_schema.read_slide_csv(input_csv)
	csv_dir = os.path.dirname(os.path.abspath(input_csv))
	if template_path:
		resolved_template, template_warnings = path_resolver.resolve_path(
			template_path,
			input_dir=csv_dir,
			strict=False,
		)
		for message in template_warnings:
			print(f"Warning: {message}")
		presentation = pptx.Presentation(resolved_template)
	else:
		presentation = pptx.Presentation()
	layout_map = build_layout_map(presentation)
	source_cache: dict[str, pptx.Presentation] = {}
	temp_dirs = []
	for row_index, row in enumerate(rows, 1):
		source_pptx = row["source_pptx"]
		source_path, path_warnings = path_resolver.resolve_source_path(
			source_pptx,
			csv_dir,
			strict=False,
		)
		for message in path_warnings:
			print(f"Warning: {message}")
		source_key = source_path
		source_presentation = source_cache.get(source_key)
		if not source_presentation:
			if source_path.lower().endswith(".odp"):
				temp_dir = tempfile.TemporaryDirectory()
				temp_dirs.append(temp_dir)
				converted = soffice_tools.convert_odp_to_pptx(source_path, temp_dir.name)
				source_presentation = pptx.Presentation(converted)
			else:
				source_presentation = pptx.Presentation(source_path)
			source_cache[source_key] = source_presentation

		slide_index = int(row["source_slide_index"])
		if slide_index < 1 or slide_index > len(source_presentation.slides):
			raise ValueError(
				f"Source slide index out of range: {source_pptx} {slide_index}."
			)
		source_slide = source_presentation.slides[slide_index - 1]
		computed_hash, _, _ = pptx_hash.compute_slide_hash_from_slide(source_slide)
		row_hash = row.get("slide_hash", "")
		if not row_hash:
			raise ValueError(f"Row {row_index}: slide_hash is missing.")
		if row_hash != computed_hash:
			raise ValueError(
				f"Row {row_index}: slide_hash mismatch for {source_pptx} slide {slide_index}."
			)
		image_blobs = collect_source_images(source_slide)
		layout_type = row.get("layout_type", "")
		if not layout_type:
			raise ValueError(f"Row {row_index}: layout_type is missing.")
		layout = select_layout(
			presentation,
			layout_map,
			row.get("master_name", ""),
			layout_type,
		)
		slide = presentation.slides.add_slide(layout)
		set_title(slide, row.get("title_text", ""))
		set_body_text(slide, row.get("body_text", ""))
		set_notes_text(slide, row.get("notes_text", ""))
		insert_images(slide, image_blobs)
	if output_path.lower().endswith(".odp"):
		with tempfile.TemporaryDirectory() as temp_dir:
			temp_pptx = os.path.join(temp_dir, "merged.pptx")
			presentation.save(temp_pptx)
			soffice_tools.convert_pptx_to_odp(temp_pptx, output_path)
	else:
		presentation.save(output_path)
	for temp_dir in temp_dirs:
		temp_dir.cleanup()
