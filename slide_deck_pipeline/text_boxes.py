# Standard Library
import hashlib
import re

# PIP3 modules
import pptx
import pptx.enum.shapes

# local repo modules
import slide_deck_pipeline.pptx_text as pptx_text


#============================================
def is_body_placeholder(placeholder_type) -> bool:
	"""
	Return True for placeholder types that act like body content.

	Args:
		placeholder_type: Placeholder type enum.

	Returns:
		bool: True if treated as body content.
	"""
	placeholders = pptx.enum.shapes.PP_PLACEHOLDER
	body_types = [placeholders.BODY]
	for attr_name in ("OBJECT", "CONTENT", "TEXT"):
		candidate = getattr(placeholders, attr_name, None)
		if candidate is not None:
			body_types.append(candidate)
	return placeholder_type in tuple(body_types)


#============================================
def normalize_shape_name(name: str) -> str:
	"""
	Normalize a shape name for stable IDs.

	Args:
		name: Shape name.

	Returns:
		str: Normalized name.
	"""
	if not name:
		return ""
	cleaned = name.strip().lower()
	cleaned = re.sub(r"[^a-z0-9]+", "_", cleaned)
	cleaned = cleaned.strip("_")
	return cleaned


#============================================
def placeholder_type_name(placeholder_type) -> str:
	"""
	Return a readable placeholder type label.

	Args:
		placeholder_type: Placeholder type enum.

	Returns:
		str: Placeholder type name.
	"""
	if not placeholder_type:
		return ""
	name = getattr(placeholder_type, "name", "")
	if not name:
		return ""
	return name.lower()


#============================================
def extract_text_block(shape) -> str:
	"""
	Extract text with indentation from a shape.

	Args:
		shape: Shape instance.

	Returns:
		str: Text block with tab indentation.
	"""
	if not getattr(shape, "has_text_frame", False):
		return ""
	lines = pptx_text.extract_paragraph_lines(shape.text_frame)
	return "\n".join(lines)


#============================================
def ensure_unique_id(box_id: str, used: set[str]) -> str:
	"""
	Ensure a box id is unique within a slide.

	Args:
		box_id: Proposed box id.
		used: Set of used ids.

	Returns:
		str: Unique box id.
	"""
	if box_id not in used:
		used.add(box_id)
		return box_id
	counter = 2
	while True:
		candidate = f"{box_id}_{counter}"
		if candidate not in used:
			used.add(candidate)
			return candidate
		counter += 1


#============================================
def collect_text_boxes(
	slide: pptx.slide.Slide,
	include_subtitle: bool,
	include_footer: bool,
	include_fallback: bool = True,
) -> tuple[list[dict[str, object]], bool]:
	"""
	Collect text boxes for export or update.

	Args:
		slide: Slide instance.
		include_subtitle: Include subtitle placeholders.
		include_footer: Include footer placeholders.
		include_fallback: Include non-placeholder shapes if needed.

	Returns:
		tuple[list[dict[str, object]], bool]: Box records and fallback flag.
	"""
	boxes = []
	used_ids: set[str] = set()
	body_count = 0
	for shape in slide.shapes:
		if not getattr(shape, "has_text_frame", False):
			continue
		if not getattr(shape, "is_placeholder", False):
			continue
		placeholder_type = shape.placeholder_format.type
		box_id = ""
		if placeholder_type in (
			pptx.enum.shapes.PP_PLACEHOLDER.TITLE,
			pptx.enum.shapes.PP_PLACEHOLDER.CENTER_TITLE,
		):
			box_id = "title"
		elif placeholder_type == pptx.enum.shapes.PP_PLACEHOLDER.SUBTITLE:
			if include_subtitle:
				box_id = "subtitle"
		elif is_body_placeholder(placeholder_type):
			body_count += 1
			box_id = f"body_{body_count}"
		elif placeholder_type == pptx.enum.shapes.PP_PLACEHOLDER.FOOTER:
			if include_footer:
				box_id = "footer"
		if not box_id:
			continue
		box_id = ensure_unique_id(box_id, used_ids)
		boxes.append(
			{
				"box_id": box_id,
				"shape": shape,
				"shape_name": getattr(shape, "name", ""),
				"placeholder_type": placeholder_type_name(placeholder_type),
			}
		)
	if boxes or not include_fallback:
		return (boxes, False)
	fallback_boxes = []
	fallback_index = 0
	for shape in slide.shapes:
		if not getattr(shape, "has_text_frame", False):
			continue
		if getattr(shape, "is_placeholder", False):
			continue
		shape_name = getattr(shape, "name", "")
		box_id = normalize_shape_name(shape_name)
		if not box_id:
			fallback_index += 1
			shape_id = getattr(shape, "shape_id", "")
			guard_source = str(shape_id or fallback_index)
			guard_hash = hashlib.sha256(guard_source.encode("utf-8")).hexdigest()[:8]
			box_id = f"box_{fallback_index}_{guard_hash}"
		box_id = ensure_unique_id(box_id, used_ids)
		fallback_boxes.append(
			{
				"box_id": box_id,
				"shape": shape,
				"shape_name": shape_name,
				"placeholder_type": "",
			}
		)
	return (fallback_boxes, True)
