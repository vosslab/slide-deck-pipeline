# Standard Library
import tempfile

# PIP3 modules
import pptx
import yaml

# local repo modules
import slide_deck_pipeline.csv_schema as csv_schema
import slide_deck_pipeline.pptx_hash as pptx_hash
import slide_deck_pipeline.pptx_io as pptx_io
import slide_deck_pipeline.pptx_text as pptx_text
import slide_deck_pipeline.text_boxes as text_boxes


#============================================
def build_box_record(shape, box_meta: dict[str, object]) -> dict[str, str]:
	"""
	Build a YAML box record from a shape.

	Args:
		shape: Shape instance.
		box_meta: Metadata for the box.

	Returns:
		dict[str, str]: Box record.
	"""
	text_value = text_boxes.extract_text_block(shape)
	box_record = {
		"box_id": box_meta["box_id"],
		"text_hash_before": csv_schema.compute_text_hash(text_value),
		"text": text_value,
	}
	shape_name = box_meta.get("shape_name", "")
	if shape_name:
		box_record["shape_name"] = shape_name
	placeholder_type = box_meta.get("placeholder_type", "")
	if placeholder_type:
		box_record["placeholder_type"] = placeholder_type
	return box_record


#============================================
def export_slide_text(
	input_path: str,
	output_path: str,
	include_notes: bool,
	include_subtitle: bool,
	include_footer: bool,
) -> None:
	"""
	Export slide text blocks to YAML.

	Args:
		input_path: Input PPTX or ODP path.
		output_path: Output YAML path.
		include_notes: Include speaker notes blocks.
		include_subtitle: Include subtitle placeholders.
		include_footer: Include footer placeholders.
	"""
	needs_conversion = input_path.lower().endswith(".odp")
	if needs_conversion:
		with tempfile.TemporaryDirectory() as temp_dir:
			pptx_path, source_name = pptx_io.resolve_input_pptx(
				input_path,
				temp_dir,
			)
			write_yaml(
				pptx_path,
				source_name,
				output_path,
				include_notes,
				include_subtitle,
				include_footer,
			)
		return
	pptx_path, source_name = pptx_io.resolve_input_pptx(input_path, None)
	write_yaml(
		pptx_path,
		source_name,
		output_path,
		include_notes,
		include_subtitle,
		include_footer,
	)


#============================================
def write_yaml(
	pptx_path: str,
	source_name: str,
	output_path: str,
	include_notes: bool,
	include_subtitle: bool,
	include_footer: bool,
) -> None:
	"""
	Write YAML from a PPTX path.

	Args:
		pptx_path: PPTX path.
		source_name: Source deck basename.
		output_path: Output YAML path.
		include_notes: Include speaker notes blocks.
		include_subtitle: Include subtitle placeholders.
		include_footer: Include footer placeholders.
	"""
	presentation = pptx.Presentation(pptx_path)
	patches = []
	fallback_slides = []
	box_count = 0
	for index, slide in enumerate(presentation.slides, 1):
		notes_text = pptx_text.extract_notes_text(slide)
		slide_hash, _, _ = pptx_hash.compute_slide_hash_from_slide(
			slide,
			notes_text,
		)
		boxes, used_fallback = text_boxes.collect_text_boxes(
			slide,
			include_subtitle,
			include_footer,
			include_fallback=True,
		)
		if used_fallback:
			fallback_slides.append(index)
		box_records = []
		for box_meta in boxes:
			shape = box_meta["shape"]
			box_records.append(build_box_record(shape, box_meta))
		if include_notes:
			box_records.append(
				{
					"box_id": "notes",
					"text_hash_before": csv_schema.compute_text_hash(notes_text),
					"text": notes_text,
					"placeholder_type": "notes",
				}
			)
		if not box_records:
			continue
		box_count += len(box_records)
		patches.append(
			{
				"source_slide_index": index,
				"slide_hash": slide_hash,
				"boxes": box_records,
			}
		)
	payload = {
		"version": 1,
		"source_pptx": source_name,
		"patches": patches,
	}
	with open(output_path, "w", encoding="utf-8") as handle:
		yaml.safe_dump(
			payload,
			handle,
			sort_keys=False,
			default_flow_style=False,
			allow_unicode=False,
		)
	print(f"Exported {len(patches)} slides with {box_count} text blocks.")
	if fallback_slides:
		listed = ", ".join(str(idx) for idx in fallback_slides)
		print(f"Fallback shape matching used on slides: {listed}")
