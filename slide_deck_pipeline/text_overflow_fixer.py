# Standard Library
import os
import tempfile

# PIP3 modules
import pptx
from pptx.enum.text import MSO_AUTO_SIZE

# local repo modules
import slide_deck_pipeline.pptx_io as pptx_io
import slide_deck_pipeline.soffice_tools as soffice_tools


#============================================
def fix_text_overflow(
	input_path: str,
	output_path: str,
	inplace: bool,
) -> tuple[int, int]:
	"""
	Enable "Shrink text on overflow" for all text boxes in a PPTX or ODP file.

	Args:
		input_path: Input PPTX or ODP path.
		output_path: Output PPTX or ODP path.
		inplace: Allow writing to the input path.

	Returns:
		tuple[int, int]: (text boxes seen, text boxes adjusted).
	"""
	if not inplace:
		input_abs = os.path.abspath(input_path)
		output_abs = os.path.abspath(output_path)
		if input_abs == output_abs:
			raise ValueError("Output path matches input; use --inplace to override.")
	output_is_odp = output_path.lower().endswith(".odp")
	if input_path.lower().endswith(".odp"):
		with tempfile.TemporaryDirectory() as temp_dir:
			pptx_path, _ = pptx_io.resolve_input_pptx(input_path, temp_dir)
			return fix_pptx(pptx_path, output_path, output_is_odp)
	pptx_path, _ = pptx_io.resolve_input_pptx(input_path, None)
	return fix_pptx(pptx_path, output_path, output_is_odp)


#============================================
def fix_pptx(
	pptx_path: str,
	output_path: str,
	output_is_odp: bool,
) -> tuple[int, int]:
	"""
	Enable "Shrink text on overflow" for all text boxes in a PPTX file.

	Args:
		pptx_path: Input PPTX path.
		output_path: Output path.
		output_is_odp: True if output should be ODP.

	Returns:
		tuple[int, int]: (text boxes seen, text boxes adjusted).
	"""
	presentation = pptx.Presentation(pptx_path)
	total = 0
	adjusted = 0
	for slide in presentation.slides:
		for shape in slide.shapes:
			if not shape.has_text_frame:
				continue
			total += 1
			text_frame = shape.text_frame
			# Check if auto_size is not already set to TEXT_TO_FIT_SHAPE
			if text_frame.auto_size != MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
				text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
				adjusted += 1
	if output_is_odp:
		with tempfile.TemporaryDirectory() as temp_dir:
			temp_pptx = os.path.join(temp_dir, "text_overflow_fixed.pptx")
			presentation.save(temp_pptx)
			soffice_tools.convert_pptx_to_odp(temp_pptx, output_path)
	else:
		presentation.save(output_path)
	return (total, adjusted)
