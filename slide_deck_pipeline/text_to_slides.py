# Standard Library
import os

# PIP3 modules
import pptx
import pptx.enum.shapes
import pptx.oxml.ns

# local repo modules
import slide_deck_pipeline.default_layouts as default_layouts
import slide_deck_pipeline.layout_classifier as layout_classifier
import slide_deck_pipeline.path_resolver as path_resolver
import slide_deck_pipeline.reporting as reporting
import slide_deck_pipeline.template as template


#============================================
def clear_text_frame(text_frame) -> None:
	"""
	Clear all paragraphs in a text frame.

	Args:
		text_frame: TextFrame instance.
	"""
	if hasattr(text_frame, "clear"):
		text_frame.clear()
	else:
		text_frame.text = ""


#============================================
def set_text(shape, text: str) -> None:
	"""
	Set text into a shape if possible.

	Args:
		shape: Shape instance.
		text: Text string.
	"""
	if not text:
		return
	if getattr(shape, "has_text_frame", False):
		shape.text_frame.text = text
		return
	if hasattr(shape, "text"):
		shape.text = text


#============================================
def set_bullets(shape, bullets: list[str]) -> None:
	"""
	Fill a body placeholder with bullets.

	Args:
		shape: Shape instance.
		bullets: Bullet strings.
	"""
	if not getattr(shape, "has_text_frame", False):
		return
	text_frame = shape.text_frame
	clear_text_frame(text_frame)
	if not bullets:
		return
	for index, bullet in enumerate(bullets):
		if index == 0:
			paragraph = text_frame.paragraphs[0]
		else:
			paragraph = text_frame.add_paragraph()
		paragraph.text = bullet
		paragraph.level = 0


#============================================
def placeholder_index(shape) -> int:
	"""
	Return a placeholder index for ordering.

	Args:
		shape: Shape instance.

	Returns:
		int: Placeholder index.
	"""
	try:
		return int(shape.placeholder_format.idx)
	except Exception:
		return 0


#============================================
def collect_placeholders(slide) -> dict[str, list[object]]:
	"""
	Collect placeholders grouped by role.

	Args:
		slide: Slide instance.

	Returns:
		dict[str, list[object]]: Placeholder lists.
	"""
	roles = {"title": [], "subtitle": [], "body": [], "picture": []}
	placeholders = getattr(slide, "placeholders", [])
	picture_types = []
	placeholder_enum = pptx.enum.shapes.PP_PLACEHOLDER
	if hasattr(placeholder_enum, "PICTURE"):
		picture_types.append(placeholder_enum.PICTURE)
	if hasattr(placeholder_enum, "MEDIA"):
		picture_types.append(placeholder_enum.MEDIA)
	for shape in placeholders:
		if not getattr(shape, "is_placeholder", False):
			continue
		try:
			placeholder_type = shape.placeholder_format.type
		except Exception:
			continue
		role = layout_classifier.classify_placeholder_role(placeholder_type)
		if role:
			roles[role].append(shape)
		if placeholder_type in picture_types:
			roles["picture"].append(shape)
	for role, items in roles.items():
		roles[role] = sorted(items, key=placeholder_index)
	return roles


#============================================
def remove_all_slides(presentation: pptx.Presentation) -> None:
	"""
	Remove all slides from a presentation.

	Args:
		presentation: Presentation to clear.
	"""
	slide_id_list = presentation.slides._sldIdLst
	partnames_to_drop: set[str] = set()
	for slide_id in list(slide_id_list):
		rel_id = slide_id.get(pptx.oxml.ns.qn("r:id"))
		slide_part = presentation.part.related_parts.get(rel_id)
		if slide_part is not None:
			partnames_to_drop.add(str(slide_part.partname))
		presentation.part.drop_rel(rel_id)
		slide_id_list.remove(slide_id)

	# Some templates include additional relationships to slide parts (for example
	# viewProps or notes slide references). Drop any remaining relationships so
	# orphaned slide parts are not marshaled back out.
	if partnames_to_drop:
		for part in presentation.part.package.iter_parts():
			for r_id, rel in list(part.rels.items()):
				if rel.is_external:
					continue
				if str(rel.target_part.partname) not in partnames_to_drop:
					continue
				try:
					part.drop_rel(r_id)
				except AttributeError:
					if r_id in part.rels:
						del part.rels[r_id]


#============================================
def resolve_images(
	entry: dict,
	input_dir: str,
	strict: bool,
	warnings: list[str],
) -> list[str]:
	"""
	Resolve image paths for a slide entry.

	Args:
		entry: Slide entry.
		input_dir: YAML input directory.
		strict: Treat warnings as errors.
		warnings: Warning list to append to.

	Returns:
		list[str]: Resolved image paths.
	"""
	images = []
	if entry.get("images"):
		images = list(entry["images"])
	elif entry.get("image"):
		images = [entry["image"]]
	resolved = []
	for item in images:
		path, path_warnings = path_resolver.resolve_path(
			item,
			input_dir=input_dir,
			strict=strict,
		)
		warnings.extend(path_warnings)
		resolved.append(path)
	return resolved


#============================================
def fill_placeholders(
	slide,
	entry: dict,
	input_dir: str,
	strict: bool,
) -> tuple[int, int, list[str]]:
	"""
	Fill placeholders on a slide.

	Args:
		slide: Slide instance.
		entry: Slide entry.
		input_dir: YAML input directory.
		strict: Treat warnings as errors.

	Returns:
		tuple[int, int, list[str]]: Images placed, images dropped, warnings.
	"""
	warnings = []
	roles = collect_placeholders(slide)
	title = entry.get("title")
	subtitle = entry.get("subtitle")
	if title:
		if roles["title"]:
			set_text(roles["title"][0], title)
		else:
			message = "Title provided but no title placeholder found."
			if strict:
				raise ValueError(message)
			warnings.append(message)
	if subtitle:
		if roles["subtitle"]:
			set_text(roles["subtitle"][0], subtitle)
		else:
			message = "Subtitle provided but no subtitle placeholder found."
			if strict:
				raise ValueError(message)
			warnings.append(message)
	bodies = entry.get("bodies", [])
	for index, placeholder in enumerate(roles["body"]):
		body = bodies[index] if index < len(bodies) else None
		if not body:
			continue
		set_bullets(placeholder, body.get("bullets", []))
	if len(bodies) > len(roles["body"]):
		message = "More bodies provided than body placeholders."
		if strict:
			raise ValueError(message)
		warnings.append(message)

	images = resolve_images(entry, input_dir, strict, warnings)
	picture_placeholders = roles["picture"]
	if images and not picture_placeholders:
		message = "Images provided but no picture placeholders found."
		if strict:
			raise ValueError(message)
		warnings.append(message)
		return (0, len(images), warnings)
	placed = 0
	dropped = 0
	for index, image_path in enumerate(images):
		if index >= len(picture_placeholders):
			dropped += 1
			continue
		placeholder = picture_placeholders[index]
		placeholder.insert_picture(image_path)
		placed += 1
	if dropped:
		message = "More images provided than picture placeholders."
		if strict:
			raise ValueError(message)
		warnings.append(message)
	return (placed, dropped, warnings)


#============================================
def render_to_pptx(
	spec: dict,
	input_path: str,
	template_override: str | None,
	output_path: str,
	strict: bool,
) -> None:
	"""
	Render a PPTX from a YAML spec.

	Args:
		spec: Normalized spec.
		input_path: Spec path.
		template_override: Template override path.
		output_path: Output PPTX path.
		strict: Treat warnings as errors.
	"""
	warnings = []
	input_dir = os.path.dirname(os.path.abspath(input_path))
	template_path = spec.get("template_deck")
	if template_override:
		template_path = template_override
	if template_path:
		resolved_template, template_warnings = path_resolver.resolve_path(
			template_path,
			input_dir=input_dir,
			strict=strict,
		)
		warnings.extend(template_warnings)
		presentation = template.load_template(resolved_template)
		remove_all_slides(presentation)
		layout_map, layout_warnings = template.build_layout_map(
			presentation,
			strict,
		)
		warnings.extend(layout_warnings)
		master_name = spec.get("defaults", {}).get("master_name", "") or ""
		def resolve_layout(layout_type: str):
			return template.resolve_layout(
				layout_map,
				master_name,
				layout_type,
				strict,
			)
	else:
		presentation = pptx.Presentation()
		layout_map, layout_warnings = default_layouts.build_default_layout_map(
			presentation,
			strict,
		)
		warnings.extend(layout_warnings)
		def resolve_layout(layout_type: str):
			return default_layouts.resolve_default_layout(
				layout_map,
				layout_type,
				strict,
			)
	total_images = 0
	total_dropped = 0
	for entry in spec.get("slides", []):
		layout_type = entry.get("layout_type")
		layout, layout_warnings = resolve_layout(layout_type)
		warnings.extend(layout_warnings)
		slide = presentation.slides.add_slide(layout)
		placed, dropped, slide_warnings = fill_placeholders(
			slide,
			entry,
			input_dir,
			strict,
		)
		total_images += placed
		total_dropped += dropped
		warnings.extend(slide_warnings)
	presentation.save(output_path)
	reporting.print_summary("Slides rendered", len(spec.get("slides", [])))
	reporting.print_summary("Images placed", total_images)
	reporting.print_summary("Images dropped", total_dropped)
	reporting.print_warnings(warnings)
