#!/usr/bin/env python3

import argparse
import os
import tempfile

# local repo modules
import slide_deck_pipeline.csv_schema as csv_schema
import slide_deck_pipeline.path_resolver as path_resolver
import slide_deck_pipeline.soffice_tools as soffice_tools


#============================================
def parse_args() -> argparse.Namespace:
	"""
	Parse command-line arguments.
	"""
	parser = argparse.ArgumentParser(
		description="Validate a merged slide CSV and related sources."
	)
	parser.add_argument(
		"-i",
		"--input",
		dest="input_csv",
		required=True,
		help="Input merged CSV path",
	)
	parser.add_argument(
		"-t",
		"--template",
		dest="template_path",
		default="",
		help="Template PPTX path for master/layout validation",
	)
	parser.add_argument(
		"-c",
		"--check-sources",
		dest="check_sources",
		help="Check source PPTX or ODP files exist",
		action="store_true",
	)
	parser.add_argument(
		"-C",
		"--no-check-sources",
		dest="check_sources",
		help="Skip source file checks",
		action="store_false",
	)
	parser.set_defaults(check_sources=True)
	parser.add_argument(
		"-s",
		"--strict",
		dest="strict",
		help="Require slide hashes to match source slides",
		action="store_true",
	)
	parser.add_argument(
		"-S",
		"--no-strict",
		dest="strict",
		help="Skip slide hash validation",
		action="store_false",
	)
	parser.set_defaults(strict=False)
	args = parser.parse_args()
	return args


#============================================
#============================================
def resolve_source_path(
	source_pptx: str,
	csv_dir: str,
	strict: bool,
	warnings: list[str],
) -> str:
	"""
	Resolve a source path using the CSV directory as fallback.

	Args:
		source_pptx: Source PPTX or ODP path.
		csv_dir: Directory containing the CSV.
		strict: Treat ambiguous matches as errors.
		warnings: Warning list to append to.

	Returns:
		str: Resolved source path.
	"""
	resolved, path_warnings = path_resolver.resolve_path(
		source_pptx,
		input_dir=csv_dir,
		strict=strict,
	)
	warnings.extend(path_warnings)
	return resolved


#============================================
def normalize_row_value(row: dict[str, str], key: str) -> str:
	"""
	Normalize a CSV row field to a string.

	Args:
		row: CSV row.
		key: Column name.

	Returns:
		str: Normalized value.
	"""
	value = row.get(key)
	if value is None:
		return ""
	if isinstance(value, str):
		return value
	return str(value)


#============================================
def normalize_name(name: str) -> str:
	"""
	Normalize a master or layout name for matching.

	Args:
		name: Name value.

	Returns:
		str: Normalized name.
	"""
	if not name:
		return ""
	return name.strip().lower().replace(" ", "_")


#============================================
def is_positive_int(value: str) -> bool:
	"""
	Check whether a string is a positive integer.

	Args:
		value: Input string.

	Returns:
		bool: True if value is a positive integer.
	"""
	if not value:
		return False
	if not value.isdigit():
		return False
	parsed = int(value)
	return parsed > 0


#============================================
def is_hex_hash(value: str) -> bool:
	"""
	Check whether a value looks like a 16-char hex hash.

	Args:
		value: Hash string.

	Returns:
		bool: True if value matches the expected format.
	"""
	if len(value) != 16:
		return False
	for ch in value.lower():
		if ch not in "0123456789abcdef":
			return False
	return True


#============================================
def load_template_layout_types(template_path: str) -> set[tuple[str, str]]:
	"""
	Load master and layout type pairs from a template PPTX.

	Args:
		template_path: Template PPTX path.

	Returns:
		set[tuple[str, str]]: Normalized (master, layout_type) pairs.
	"""
	# PIP3 modules
	import pptx

	# local repo modules
	import slide_deck_pipeline.layout_classifier as layout_classifier

	presentation = pptx.Presentation(template_path)
	slide_width = int(getattr(presentation, "slide_width", 0) or 0)
	slide_height = int(getattr(presentation, "slide_height", 0) or 0)
	available = set()
	for layout in presentation.slide_layouts:
		layout_type, _, _ = layout_classifier.classify_layout_type(
			layout,
			slide_width,
			slide_height,
			"",
			"",
		)
		master = getattr(layout, "slide_master", None)
		master_name = normalize_name(getattr(master, "name", ""))
		if not layout_type:
			continue
		available.add((master_name, layout_type))
	return available


#============================================
def validate_rows(
	rows: list[dict[str, str]],
	csv_dir: str,
	check_sources: bool,
	strict: bool,
	template_path: str,
) -> tuple[list[str], list[str]]:
	"""
	Validate merged CSV rows.

	Args:
		rows: CSV rows.
		csv_dir: Directory containing the CSV.
		check_sources: Whether to check source files exist.
		strict: Whether to validate slide hashes against sources.
		template_path: Template PPTX path for layout validation.

	Returns:
		tuple[list[str], list[str]]: Errors and warnings.
	"""
	errors = []
	warnings = []
	layout_pairs: set[tuple[str, str]] = set()
	if template_path:
		try:
			resolved_template, template_warnings = path_resolver.resolve_path(
				template_path,
				input_dir=csv_dir,
				strict=strict,
			)
			warnings.extend(template_warnings)
			layout_pairs = load_template_layout_types(resolved_template)
		except FileNotFoundError:
			errors.append("Template PPTX not found.")
	if not rows:
		warnings.append("No rows found in CSV.")
		return (errors, warnings)

	source_cache: dict[str, object] = {}
	temp_dirs: list[tempfile.TemporaryDirectory] = []
	pptx_module = None
	pptx_hash_module = None
	if strict:
		# PIP3 modules
		import pptx as pptx_module

		# local repo modules
		import slide_deck_pipeline.pptx_hash as pptx_hash_module
	for index, row in enumerate(rows, 1):
		source_pptx = normalize_row_value(row, "source_pptx")
		if not source_pptx:
			errors.append(f"Row {index}: missing source_pptx.")
		else:
			extension = os.path.splitext(source_pptx)[1].lower()
			if extension not in (".pptx", ".odp"):
				warnings.append(f"Row {index}: unexpected source_pptx extension.")
			if check_sources or strict:
				try:
					resolved_path = resolve_source_path(
						source_pptx,
						csv_dir,
						strict,
						warnings,
					)
				except FileNotFoundError:
					errors.append(f"Row {index}: source_pptx not found.")
					resolved_path = ""
				if resolved_path and not os.path.exists(resolved_path):
					errors.append(f"Row {index}: source_pptx not found.")

		slide_index = normalize_row_value(row, "source_slide_index")
		if not is_positive_int(slide_index):
			errors.append(f"Row {index}: invalid source_slide_index {slide_index}.")

		slide_hash = normalize_row_value(row, "slide_hash")
		if not slide_hash:
			errors.append(f"Row {index}: missing slide_hash.")
		elif not is_hex_hash(slide_hash):
			errors.append(f"Row {index}: slide_hash must be 16 hex characters.")

		master_name = normalize_row_value(row, "master_name")
		layout_type = normalize_row_value(row, "layout_type")
		layout_type_key = normalize_name(layout_type)
		if not master_name:
			errors.append(f"Row {index}: missing master_name.")
		if not layout_type:
			errors.append(f"Row {index}: missing layout_type.")
		if layout_pairs and master_name and layout_type:
			pair = (normalize_name(master_name), layout_type_key)
			if pair not in layout_pairs:
				errors.append(f"Row {index}: master/layout_type not found in template.")

		if strict and source_pptx and is_positive_int(slide_index) and slide_hash:
			try:
				resolved_path = resolve_source_path(
					source_pptx,
					csv_dir,
					strict,
					warnings,
				)
			except FileNotFoundError:
				continue
			source_presentation = source_cache.get(resolved_path)
			if not source_presentation:
				if resolved_path.lower().endswith(".odp"):
					temp_dir = tempfile.TemporaryDirectory()
					temp_dirs.append(temp_dir)
					converted = soffice_tools.convert_odp_to_pptx(resolved_path, temp_dir.name)
					source_presentation = pptx_module.Presentation(converted)
				else:
					source_presentation = pptx_module.Presentation(resolved_path)
				source_cache[resolved_path] = source_presentation
			slide_number = int(slide_index)
			if slide_number < 1 or slide_number > len(source_presentation.slides):
				errors.append(
					f"Row {index}: source_slide_index out of range for {source_pptx}."
				)
				continue
			source_slide = source_presentation.slides[slide_number - 1]
			computed_hash, _, _ = pptx_hash_module.compute_slide_hash_from_slide(
				source_slide
			)
			if computed_hash != slide_hash:
				errors.append(f"Row {index}: slide_hash mismatch.")

	for temp_dir in temp_dirs:
		temp_dir.cleanup()
	return (errors, warnings)


#============================================
def format_messages(label: str, messages: list[str]) -> list[str]:
	"""
	Format validation messages with a label.

	Args:
		label: Message label.
		messages: List of messages.

	Returns:
		list[str]: Formatted lines.
	"""
	lines = []
	for message in messages:
		lines.append(f"{label}: {message}")
	return lines


#============================================
def main() -> None:
	"""
	Main entry point.
	"""
	args = parse_args()
	rows = csv_schema.read_slide_csv(args.input_csv)
	csv_dir = os.path.dirname(os.path.abspath(args.input_csv))
	errors, warnings = validate_rows(
		rows,
		csv_dir,
		args.check_sources,
		args.strict,
		args.template_path,
	)
	if warnings:
		for line in format_messages("WARN", warnings):
			print(line)
	if errors:
		for line in format_messages("ERROR", errors):
			print(line)
		raise RuntimeError(f"CSV validation failed with {len(errors)} errors.")
	print("CSV validation OK.")


if __name__ == "__main__":
	main()
